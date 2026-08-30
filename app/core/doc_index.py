"""Document index — text extraction, chunking, and per-project index persistence.

Roadmap V2 · Stream C — Phase C1.

This module builds a searchable text index for each project's documents.
It reuses file_crypto for transparent decryption and app.core.projects for
document metadata. No TF-IDF or search logic lives here — that is Phase C2.

Public API
----------
* ``extract_document_text(file_path, filename) -> str``
    Extract plaintext from a single document. Returns "" on any error.
* ``chunk_text(text, words_per_chunk=500) -> list[str]``
    Split text into word-count-bounded chunks.
* ``index_project(project_id) -> dict``
    Build (or rebuild) the full index for a project and persist it.
* ``index_document(project_id, document_id) -> dict``
    Incrementally add / replace one document in an existing index.
* ``invalidate_project(project_id) -> None``
    Delete the on-disk index so the next call rebuilds from scratch.
* ``_load_index(project_id) -> dict | None``
    Read a project's stored index; None if absent.
* ``init_db() -> None``
    Ensure the index DB schema exists and import any legacy JSON indexes.
* ``_data_dir() -> str``
    DATA_DIR env at call time, with tempfile fallback.

The index is persisted in the unified The Fork DB (one JSON row per project),
so the read-modify-write in ``index_document`` runs in a real transaction —
concurrent updates, including across worker processes, serialise instead of
overwriting each other.

SQLAlchemy-backed via app.core.db — unified The Fork schema.
"""

from __future__ import annotations

import asyncio
import concurrent.futures
import copy
import json
import logging
import logging as _logging
import os
import re
import sqlite3
import tempfile
import threading
import zlib
from collections.abc import Callable
from datetime import datetime, timezone
from typing import Any

from sqlalchemy import delete, insert, select, text, update
from sqlalchemy.engine import Connection
from sqlalchemy.orm import Session
from sqlalchemy.orm.attributes import flag_modified

from app.core import file_crypto
from app.core import projects as _projects
from app.core.db import SessionLocal, engine, get_database_url, get_engine
from app.core.models import DocIndex, Project
from app.core.subprocess_env import scrubbed_env

logger = logging.getLogger(__name__)

# Image extensions — Stream F runs OCR on these to make scanned drawings /
# photos searchable. They are SUPPORTED (not "unsupported_type") even when OCR
# yields no text: a blank photo simply indexes with empty chunks.
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tif", ".tiff"}

# Extensions we know how to extract text from.
_SUPPORTED_EXTS = (
    {".txt", ".md", ".csv", ".json", ".xml", ".pdf", ".doc", ".docx", ".xlsx", ".pptx", ".kmz", ".zip", ".rar", ".msg", ".ifc"}
    | _IMAGE_EXTS
)

# A PDF whose recovered text-layer is shorter than this is treated as a
# scanned / image-only PDF and re-extracted via OCR.
_PDF_OCR_THRESHOLD = 30

# Skip PDFs above this size entirely. Multi-hundred-MB scanned drawing sets
# (e.g. 450 MB) OOM the 2 GB Render worker during fitz load / OCR.
_DEFAULT_PDF_MAX_SIZE_MB = 100.0
# For PDFs above this size, do NOT run OCR — only extract any existing text
# layer. 25 MB was too tight: the live priced BOQ (doc 20ac033d,
# IP-INF-053-…-BOQ-… (Priced).pdf) is ~28 MB and is a scan, so item codes
# (D599.5, D549.2) never entered the search index. 32 MB covers that file
# with headroom; page-at-a-time OCR + isolated children still bound memory.
_DEFAULT_PDF_OCR_MAX_SIZE_MB = 32.0
# Default OCR page budget. 40 left later demolition items unindexed
# ("PARTIAL searchable"). 80 covers a typical bill; BOQ-named files get a
# higher cap via PDF_OCR_BOQ_PAGE_CAP.
_DEFAULT_PDF_OCR_PAGE_CAP = 80
_DEFAULT_PDF_OCR_BOQ_PAGE_CAP = 160
# Admin / force-OCR budget. The live priced BOQ (20ac033d) is 370 pages;
# 160 left later CESMM rows (D599.5 / D549.2) unindexed even when OCR ran.
_DEFAULT_PDF_OCR_FORCE_PAGE_CAP = 400
# Isolated page-batch size for large scans. A timeout/OOM on pages 41-60
# then keeps pages 1-40 instead of ZERO_CHUNK'ing the whole document.
_DEFAULT_PDF_OCR_BATCH_PAGES = 20
# Files above this size extract in isolated page batches (the old 20 MB
# BOQ parser / OCR-skip cliff).
_DEFAULT_PDF_BATCH_MIN_MB = 20.0

# Ceiling on the TEXT a single document may accumulate. The size guards above
# bound the input; nothing bounded the output, so a long text-layer PDF grew
# `parts` without limit and paid for it twice more in chunking and embedding.
# Set well above anything real in this corpus (the largest Vol 2 specification
# extracts ~2.2M chars) so it only catches runaway documents.
_MAX_EXTRACT_CHARS = int(os.getenv("DOC_MAX_EXTRACT_CHARS", "12000000"))


def _env_float(name: str, default: float) -> float:
    raw = (os.getenv(name) or "").strip()
    if not raw:
        return default
    try:
        value = float(raw)
    except ValueError:
        return default
    return value if value >= 0 else default


def _env_int(name: str, default: int) -> int:
    raw = (os.getenv(name) or "").strip()
    if not raw:
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    return value if value >= 0 else default


def pdf_max_size_mb() -> float:
    """Skip-the-whole-PDF ceiling. Read at call time so env changes apply."""
    return _env_float("PDF_MAX_SIZE_MB", _DEFAULT_PDF_MAX_SIZE_MB)


def pdf_ocr_max_size_mb(filename: str = "") -> float:
    """Above this, OCR is skipped (text layer only). Call-time env read.

    BOQ-named files use ``PDF_OCR_BOQ_MAX_SIZE_MB`` (default: the whole-PDF
    ceiling). A leftover dashboard ``PDF_OCR_MAX_SIZE_MB=25`` must not
    re-skip the live priced bill, and Fernet ciphertext (~33% larger than
    the recorded plaintext size) must not be compared against this gate —
    callers must pass ``_document_size_mb``.
    """
    if filename and _BOQ_NAME_RE.search(filename):
        return _env_float("PDF_OCR_BOQ_MAX_SIZE_MB", pdf_max_size_mb())
    return _env_float("PDF_OCR_MAX_SIZE_MB", _DEFAULT_PDF_OCR_MAX_SIZE_MB)


def pdf_ocr_page_cap(filename: str = "") -> int:
    """OCR page budget for this PDF.

    BOQ-named files (priced bills, demolition BOQs) use a higher default so
    later item codes are not dropped by the generic cap. Read at call time
    so a dashboard env change takes effect without re-importing this module.
    """
    if filename and _BOQ_NAME_RE.search(filename):
        return _env_int("PDF_OCR_BOQ_PAGE_CAP", _DEFAULT_PDF_OCR_BOQ_PAGE_CAP)
    return _env_int("PDF_OCR_PAGE_CAP", _DEFAULT_PDF_OCR_PAGE_CAP)


def pdf_ocr_force_page_cap() -> int:
    """Admin reindex budget so a 370-page scanned BOQ is not capped at 160."""
    return _env_int("PDF_OCR_FORCE_PAGE_CAP", _DEFAULT_PDF_OCR_FORCE_PAGE_CAP)


def _document_size_mb(file_path: str) -> float:
    """Plaintext megabytes. ``os.path.getsize`` is ciphertext when encryption
    is on (Fernet base64 inflates ~33%) and is the wrong input to the OCR
    skip — 26.9 MB recorded / ~36 MB on disk sat above the 32 MB gate.
    """
    try:
        return file_crypto.plaintext_size(file_path) / (1024 * 1024)
    except Exception:
        logger.warning(
            "plaintext_size failed for %s; falling back to on-disk size",
            file_path, exc_info=True,
        )
        try:
            return os.path.getsize(file_path) / (1024 * 1024)
        except OSError:
            return 0.0


def pdf_ocr_batch_pages() -> int:
    return max(1, _env_int("PDF_OCR_BATCH_PAGES", _DEFAULT_PDF_OCR_BATCH_PAGES))


def pdf_batch_min_mb() -> float:
    return _env_float("PDF_BATCH_MIN_MB", _DEFAULT_PDF_BATCH_MIN_MB)

# In-process guard around index writes. Cross-process safety comes from the
# SQLite BEGIN IMMEDIATE transaction in _update_index; this lock just avoids
# threads in one process contending on the DB lock unnecessarily.
_INDEX_LOCK = threading.RLock()
_initialized = False
_initialized_for_url: str | None = None


# ── sync → async bridge ────────────────────────────────────────────────────────

def _run_sync(coro):
    """Run an async coroutine to completion from synchronous code.

    ``extract_document_text`` is sync but the OCR block is async. This bridge
    handles both call contexts:

    * No event loop running (plain sync indexing) → ``asyncio.run``.
    * A loop IS already running (``search_project_documents`` lazily triggers
      ``index_project`` → ``extract_document_text`` from inside the running
      loop) → run the coroutine on a FRESH loop inside a worker thread, so we
      never hit "asyncio.run() cannot be called from a running event loop".
    """
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        # No loop running in this thread — safe to drive one directly.
        return asyncio.run(coro)

    # A loop is already running here; offload to a worker thread with its own.
    def _worker():
        loop = asyncio.new_event_loop()
        try:
            return loop.run_until_complete(coro)
        finally:
            loop.close()

    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
        return pool.submit(_worker).result()


def _ocr_extract(file_path: str) -> tuple[str, bool]:
    """Run the OCR block on ``file_path``; return ``(text, low_quality)``.

    Never raises — any OCR error / missing text yields ``("", False)``. The
    OCR block does its own ``open_plaintext`` decryption, so the raw stored
    path is passed straight through.
    """
    try:
        from app.blocks.ocr import OCRBlock

        result = _run_sync(OCRBlock().process(file_path))
        if not isinstance(result, dict) or result.get("status") == "error":
            return "", False
        text = result.get("text") or ""
        quality = result.get("quality") or {}
        low_quality = bool(quality.get("low_quality"))
        return text, low_quality
    except Exception:
        return "", False


def _safety_world_extract(file_path: str, filename: str) -> str:
    """Run the baked YOLO-Worldv2 detector and return a searchable summary.

    The detector is only loaded when ``SAFETY_WORLD_WEIGHTS`` points at an
    existing baked .onnx file. Detections are summarised as a plain-text
    sentence so construction photos that OCR blanks still produce RAG
    chunks. Never raises — any failure yields "".
    """
    try:
        from pathlib import Path

        from app.blocks.safety_world_detector import default_detector
        from app.core.file_crypto import open_plaintext

        detector = default_detector()
        if detector is None:
            return ""

        conf = float(os.getenv("SAFETY_WORLD_CONF", "0.05"))
        with open_plaintext(file_path) as plain_path:
            detections = detector.detect(Path(plain_path), conf_threshold=conf)
        if not detections:
            return ""

        # Group by class, keep highest confidence per class, stable order.
        by_class: dict[str, float] = {}
        for d in detections:
            cls = d.get("class", "unknown")
            conf_val = float(d.get("confidence") or 0.0)
            by_class[cls] = max(by_class.get(cls, 0.0), conf_val)

        items = ", ".join(
            f"{cls} ({conf_val:.2f})"
            for cls, conf_val in sorted(by_class.items(), key=lambda kv: (-kv[1], kv[0]))
        )
        return f"Construction site photo {filename}: detected {items}."
    except Exception:
        return ""


def _ocr_dpi_for_page(width_pt: float, height_pt: float,
                      base_dpi: int, max_pixels: int) -> int:
    """The dpi to render a page at so the bitmap never exceeds ``max_pixels``.

    F26 (live OOM 2026-08-15): dpi alone does not bound memory — an A1 CAD
    sheet at 150 dpi is a ~35-megapixel bitmap, and rendering one during a
    sync reindex took the box from 585 MB to the 2 GB ceiling in one minute
    and dropped the service. The PIXEL COUNT is the real bound: small pages
    keep the full dpi, large sheets scale down proportionally (floor 50 dpi
    so text stays legible enough for OCR to try).
    """
    w_in = max(float(width_pt), 1.0) / 72.0
    h_in = max(float(height_pt), 1.0) / 72.0
    px = (w_in * base_dpi) * (h_in * base_dpi)
    if px <= max_pixels:
        return base_dpi
    return max(50, int(base_dpi * (max_pixels / px) ** 0.5))


def _ocr_pdf_page(page) -> str:
    """OCR a single ``fitz`` page via a rendered pixmap.

    MEMORY-BOUNDED twice over: the rendered pixel count is capped via
    ``_ocr_dpi_for_page`` (large CAD sheets scale down instead of ballooning
    — see F26), and the pixmap is dropped and the temp PNG deleted before
    returning, so OCR'ing a long scan one page at a time can never accumulate
    bitmaps and OOM the 2 GB box. Never raises — returns "" on any error.
    """
    try:
        import pytesseract
        from PIL import Image

        from app.blocks.ocr import _ocr_lang
    except Exception:
        return ""
    base_dpi = int(os.getenv("PDF_OCR_DPI", "150"))
    max_px = int(os.getenv("PDF_OCR_MAX_PIXELS", "6000000"))
    tmp_path: str | None = None
    pix = None
    try:
        # Inside the try: an exotic page object without .rect must degrade to
        # "" like every other OCR failure, never raise out of the helper.
        dpi = _ocr_dpi_for_page(page.rect.width, page.rect.height, base_dpi, max_px)
        pix = page.get_pixmap(dpi=dpi)
        fd, tmp_path = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        pix.save(tmp_path)
        with Image.open(tmp_path) as img:
            return pytesseract.image_to_string(img, lang=_ocr_lang())
    except Exception:
        return ""
    finally:
        pix = None  # free the bitmap immediately
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except Exception:
                logger.warning(
                    "swallowed %s in _ocr_pdf_page() — continuing",
                    "Exception", exc_info=True,
                )


def _pdf_tables_enabled(file_path: str, readable_path: str | None = None) -> bool:
    """Whether to run pdfplumber table extraction on this PDF.

    pdfplumber loads the whole PDF into memory, which is the OOM hazard on the
    2 GB box for large scans. Skip it above ``PDF_TABLES_MAX_MB`` — per-page OCR
    still captures the text, and image-only scans have no extractable tables
    anyway. Digital table BOQs (small files) keep the table extraction.

    F26b (live OOM 2026-08-15, the SECOND one): file size is not the only
    proxy for cost. A 5 MB vector CAD plot carried ~110,000 drawing objects
    PER A1 PAGE; pdfplumber materialises a Python object for every path, and
    the sync reindex took the box down again after the pixmap bound shipped.
    Sheets larger than ``PDF_TABLES_MAX_PAGE_PT`` on their long side (default
    1250 pt ≈ just over A3) are drawings, not table BOQs — table extraction
    is skipped for them outright.

    THIRD instance of the same lesson (2026-08-23): neither gate above catches
    a long ORDINARY document. Vol 2 Specification parts 4 and 8 are plain A4,
    7.5 MB and 12.8 MB — under the size bar, under the page-size bar — but 738
    and 387 pages. Measured on those two files, table extraction is essentially
    the entire cost::

        part 4 (738 pp)   tables on: 4191 MB / 173s     off: 122 MB / 10s
        part 8 (387 pp)   tables on: 2094 MB /  93s     off: 154 MB /  6s

    That is ~97% of peak memory and ~94% of wall-clock, and 4191 MB is what
    OOM-killed a 4 GB instance outright before extraction was isolated. It
    also explains a document that SUCCEEDED: a 36 MB / 419-page file extracted
    in 3s because the size gate had already skipped pdfplumber for it.

    Cost scales with page count at roughly 5.5 MB/page on this corpus, so the
    default of 200 pages holds pdfplumber near 1.1 GB — inside the extraction
    child's 1536 MB budget with room for the text layer (~150 MB) and OCR.
    Losing table extraction on such a document costs structure, not content:
    the cell text is already in the text layer that fitz reads. Digital table
    BOQs, which are what this feature exists for, are tens of pages and are
    unaffected.
    """
    try:
        # Default matches the 32 MB OCR plaintext gate so a 26.9 MB priced
        # BOQ is not skipped on size alone. Measure plaintext — Fernet
        # ciphertext is ~33% larger (same class as #449). An explicit
        # PDF_TABLES_MAX_MB still wins (ops OOM knob). The page-count
        # gate below still skips a 368-page scan.
        raw = (os.getenv("PDF_TABLES_MAX_MB") or "").strip()
        if raw:
            try:
                max_mb = float(raw)
            except ValueError:
                max_mb = _DEFAULT_PDF_OCR_MAX_SIZE_MB
        else:
            max_mb = max(pdf_ocr_max_size_mb(), _DEFAULT_PDF_OCR_MAX_SIZE_MB)
        if max_mb > 0 and _document_size_mb(file_path) > max_mb:
            return False
        max_page_pt = float(os.getenv("PDF_TABLES_MAX_PAGE_PT", "1250"))
        max_pages = int(os.getenv("PDF_TABLES_MAX_PAGES", "200"))
        probe = readable_path or file_path
        try:
            import fitz
            with fitz.open(probe) as _doc:
                if max_pages > 0 and _doc.page_count > max_pages:
                    return False
                for _page in _doc:
                    if max(_page.rect.width, _page.rect.height) > max_page_pt:
                        return False
        except Exception:
            # Unreadable via fitz here — let the main extractor decide; the
            # size gate above still holds. LOGGED, not silent: a probe that
            # always throws would silently disable the page-size half of the
            # OOM guard, which is precisely how the second dense-CAD OOM got
            # through (F26b).
            logger.warning(
                "swallowed %s probing page size in _pdf_tables_enabled(%s) — "
                "size gate still applies",
                "Exception", probe, exc_info=True,
            )
        return True
    except Exception:
        return False


def _pdf_tables_markdown(plumber_page) -> str:
    """Extract a pdfplumber page's tables as pipe-delimited rows so BOQ /
    schedule line items survive as retrievable text. Never raises."""
    try:
        tables = plumber_page.extract_tables() or []
    except Exception:
        return ""
    rows: list[str] = []
    for tbl in tables:
        for row in tbl or []:
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
    return "\n".join(rows)


_OOM_MARKERS = (
    "out of memory",
    "cannot allocate",
    "bad_alloc",
    "insufficient memory",
)
_OOM_ALLOC_RE = re.compile(r"\b(?:c|m|re)alloc\b[^:]*fail", re.IGNORECASE)


def _is_memory_exhaustion(exc: BaseException) -> bool:
    """True when ``exc`` -- or anything it wraps -- is an out-of-memory failure.

    ``RLIMIT_AS`` does NOT reliably surface as :class:`MemoryError`. Only
    allocations made through CPython's allocator raise it; native extractors
    call their own C allocators and raise their own types. Observed live on
    2026-08-23, inside the extraction child, on one document::

        RuntimeError: code=2: calloc (4104 x 1 bytes) failed   <- PyMuPDF
        MemoryError                                            <- pdfminer
        pdfplumber.utils.exceptions.PdfminerException          <- the wrapper

    Only the middle one is a ``MemoryError``, and it was re-raised inside a
    ``close()`` during cleanup, so the exception that actually reached the
    handler was a plain ``RuntimeError``. ``except MemoryError: raise`` never
    matched, ``except Exception: return "", {}`` did, and a document that had
    already yielded hundreds of good pages was reported as a *successful empty
    file* -- which then indexed as ZERO_CHUNK.

    Walking ``__cause__``/``__context__`` catches the wrapped case; the text
    markers catch the native case, where there is no ``MemoryError`` anywhere
    in the chain to find.
    """
    seen: set[int] = set()
    cur: BaseException | None = exc
    while cur is not None and id(cur) not in seen:
        seen.add(id(cur))
        if isinstance(cur, MemoryError):
            return True
        text = f"{type(cur).__name__}: {cur}"
        if _OOM_ALLOC_RE.search(text):
            return True
        if any(marker in text.lower() for marker in _OOM_MARKERS):
            return True
        cur = cur.__cause__ or cur.__context__
    return False


def _extract_pdf(
    file_path: str,
    filename: str = "",
    *,
    start_page: int = 0,
    end_page: int | None = None,
    ocr_budget: int | None = None,
    force_ocr: bool = False,
) -> tuple[str, dict[str, Any]]:
    """Per-page PDF extraction.

    For each page: keep its text layer when present, else OCR THAT page. Also
    append pdfplumber table rows. This fixes the cover-page bug where a digital
    cover page (> ``_PDF_OCR_THRESHOLD`` chars) made the WHOLE PDF skip OCR even
    though the body pages were image-only.

    MEMORY-BOUNDED for the 2 GB box: OCR at most ``pdf_ocr_page_cap`` pages
    (sets ``ocr_truncated``), freeing each pixmap + temp file per page. Never
    raises — returns ``("", {})`` on any error.

    ``start_page`` / ``end_page`` / ``ocr_budget`` are the isolated page-batch
    contract: a parent can OCR pages 20-40 with a remaining budget without
    re-applying the whole-file size gate (the parent already decided).

    ``force_ocr`` (admin reindex) ignores the OCR size skip whenever the
    plaintext file is under the whole-PDF ceiling, and uses the force page
    budget so a 370-page scanned BOQ is not capped at 160.
    """
    import fitz

    parts: list[str] = []
    ocr_pages = 0
    ocr_attempts = 0
    empty_text_pages = 0
    truncated = False
    chars = 0
    text_truncated = False
    ocr_skipped = False
    try:
        if ocr_budget is None:
            page_cap = pdf_ocr_page_cap(filename)
            size_mb = _document_size_mb(file_path)
            max_mb = pdf_max_size_mb()
            if max_mb > 0 and size_mb > max_mb:
                return "", {"skipped_too_large": True, "size_mb": round(size_mb, 1)}
            # Large scans: disable OCR entirely, rely on text layer only —
            # unless admin reindex asked to OCR anyway (plaintext still under
            # the hard ceiling).
            ocr_max = pdf_ocr_max_size_mb(filename)
            if force_ocr:
                ocr_max = max(ocr_max, max_mb if max_mb > 0 else ocr_max)
                page_cap = max(page_cap, pdf_ocr_force_page_cap())
            if ocr_max > 0 and size_mb > ocr_max:
                page_cap = 0
                ocr_skipped = True
        else:
            page_cap = max(0, int(ocr_budget))
        with file_crypto.open_plaintext(file_path) as readable_path:
            plumber = None
            # Skip pdfplumber on large scans — it loads the whole PDF and is
            # the OOM hazard; per-page OCR still captures the text.
            if _pdf_tables_enabled(file_path, readable_path):
                try:
                    import pdfplumber
                    plumber = pdfplumber.open(readable_path)
                except Exception:
                    plumber = None
            doc = fitz.open(readable_path)
            try:
                for i, page in enumerate(doc):
                    if i < start_page:
                        continue
                    if end_page is not None and i >= end_page:
                        break
                    # Stop accumulating past the text ceiling. Breaking keeps
                    # what was read rather than discarding the document, and
                    # the flag below tells the caller the tail is missing
                    # instead of letting a partial doc look complete.
                    if _MAX_EXTRACT_CHARS > 0 and chars >= _MAX_EXTRACT_CHARS:
                        text_truncated = True
                        break
                    page_text = (page.get_text() or "").strip()
                    # Always keep any real text layer — never discard digital
                    # text in favour of OCR.
                    if page_text:
                        parts.append(page_text)
                        chars += len(page_text)
                    # A page whose text layer is too thin is image-only (or
                    # near-empty) — OCR it, bounded by the page cap so a long
                    # scan can't OOM the box.
                    if len(page_text) < _PDF_OCR_THRESHOLD:
                        empty_text_pages += 1
                        if page_cap > 0 and ocr_attempts < page_cap:
                            ocr_attempts += 1
                            ocr_text = _ocr_pdf_page(page)
                            if ocr_text.strip():
                                parts.append(ocr_text)
                                chars += len(ocr_text)
                                ocr_pages += 1
                        elif page_cap > 0:
                            truncated = True
                    if plumber is not None and i < len(plumber.pages):
                        table_md = _pdf_tables_markdown(plumber.pages[i])
                        if table_md:
                            parts.append(table_md)
                            chars += len(table_md)
            finally:
                doc.close()
                if plumber is not None:
                    try:
                        plumber.close()
                    except Exception:
                        logger.warning(
                            "swallowed %s in _extract_pdf() — continuing",
                            "Exception", exc_info=True,
                        )
    except MemoryError:
        # NEVER report a memory failure as an empty document. MemoryError is a
        # subclass of Exception, so the handler below used to absorb it and
        # return ("", {}) -- indistinguishable from a genuinely empty file.
        #
        # That is exactly how the isolation rollout broke production silently:
        # the child hit its RLIMIT_AS ceiling, this handler turned it into a
        # clean empty extraction, and every document >= 1 MB indexed as
        # ZERO_CHUNK with no error logged anywhere. Let it propagate so the
        # caller can say what actually happened.
        raise
    except Exception as exc:
        if _is_memory_exhaustion(exc):
            # The child hit its ceiling. Everything already in `parts` was
            # extracted successfully -- discarding it turns a document that is
            # 90% readable into a ZERO_CHUNK failure, which is how a 738-page
            # spec indexed as empty while its first 40 pages sat in memory.
            #
            # Keep the pages we have and SAY the tail is missing. Callers can
            # tell a partial document from a complete one; nothing downstream
            # gets to mistake this for a genuinely empty file.
            if parts:
                logger.warning(
                    "PDF extraction ran out of memory after %d pages (%d chars) "
                    "— keeping the partial text: %s",
                    len(parts), chars, exc,
                )
                partial_meta: dict[str, Any] = {
                    "extract_oom_partial": True,
                    "extract_chars": chars,
                    "extract_parts": len(parts),
                }
                if ocr_pages > 0:
                    partial_meta["ocr_low_quality"] = True
                    partial_meta["ocr_pages"] = ocr_pages
                if ocr_attempts:
                    partial_meta["ocr_attempts"] = ocr_attempts
                if empty_text_pages:
                    partial_meta["empty_text_pages"] = empty_text_pages
                if empty_text_pages > 0 and ocr_attempts == 0:
                    partial_meta["ocr_required"] = True
                return "\n".join(parts), partial_meta
            # Nothing salvageable. Surface it as the memory failure it is
            # rather than as an empty document -- the ambiguity that hid this
            # bug for two days.
            raise MemoryError(f"PDF extraction exhausted memory: {exc}") from exc
        return "", {}
    meta: dict[str, Any] = {}
    if ocr_pages > 0:
        # The doc relied on OCR for some pages — flag for lower confidence.
        meta["ocr_low_quality"] = True
        meta["ocr_pages"] = ocr_pages
    if ocr_attempts:
        meta["ocr_attempts"] = ocr_attempts
    if empty_text_pages:
        meta["empty_text_pages"] = empty_text_pages
    if empty_text_pages > 0 and ocr_attempts == 0:
        # Cover-only extract: finer chunker can still emit 6 chunks and look
        # like success. Callers must not report that as an indexed body.
        # A single blank digital page that we *did* OCR (attempted, empty
        # result) is not this failure — that is a real blank, not a skip.
        meta["ocr_required"] = True
    if truncated:
        meta["ocr_truncated"] = True
    if text_truncated:
        meta["text_truncated"] = True
        meta["extract_chars"] = chars
    if ocr_skipped:
        meta["ocr_skipped_too_large"] = True
    return "\n".join(parts), meta


def _pdf_should_batch(file_path: str) -> bool:
    """True when this PDF should extract in isolated page batches.

    The old 20 MB cliff (BOQ parser refuse + OCR skip at 25 MB) is exactly
    the live priced-BOQ size band. Batching keeps each isolated child to
    ``PDF_OCR_BATCH_PAGES`` so a timeout/OOM on a later range does not
    ZERO_CHUNK pages already read.
    """
    size_mb = _document_size_mb(file_path)
    return size_mb > pdf_batch_min_mb()


def _pdf_page_count(file_path: str) -> int:
    """Page count for batching. 0 when the file cannot be opened."""
    try:
        import fitz

        with file_crypto.open_plaintext(file_path) as readable_path:
            doc = fitz.open(readable_path)
            try:
                return int(doc.page_count)
            finally:
                doc.close()
    except Exception:
        logger.warning(
            "could not read PDF page count for %s", file_path, exc_info=True
        )
        return 0


def _extract_pdf_range_job(
    file_path: str,
    filename: str,
    start: int,
    end: int,
    ocr_budget: int,
) -> tuple[str, dict[str, Any]]:
    """Picklable entry point for an isolated page-range extraction."""
    return _extract_pdf(
        file_path,
        filename,
        start_page=start,
        end_page=end,
        ocr_budget=ocr_budget,
    )


def _extract_pdf_batched(
    file_path: str, filename: str, *, force_ocr: bool = False
) -> tuple[str, dict[str, Any]]:
    """Extract a large PDF as isolated page batches, keeping prior text.

    Each batch runs in its own RLIMIT_AS child. A timeout or memory failure
    on batch N keeps batches 1..N-1 and flags ``extract_partial`` instead of
    returning an empty document (the isolation fallback that used to wipe a
    28 MB scan that had already OCR'd its first pages).
    """
    from app.core.extract_isolated import run_isolated

    size_mb = _document_size_mb(file_path)
    max_mb = pdf_max_size_mb()
    if max_mb > 0 and size_mb > max_mb:
        return "", {"skipped_too_large": True, "size_mb": round(size_mb, 1)}

    n_pages = _pdf_page_count(file_path)
    if n_pages <= 0:
        (text, meta), diag = run_isolated(
            _extract_pdf,
            (file_path, filename),
            fallback=("", {}),
            label=f"extraction of {filename}",
        )
        if diag:
            meta = {**(meta or {}), **diag}
        return text or "", meta or {}

    ocr_max = pdf_ocr_max_size_mb(filename)
    if force_ocr:
        ocr_max = max(ocr_max, max_mb if max_mb > 0 else ocr_max)
    ocr_skipped = ocr_max > 0 and size_mb > ocr_max
    if ocr_skipped:
        budget = 0
    elif force_ocr:
        budget = max(pdf_ocr_page_cap(filename), pdf_ocr_force_page_cap(), n_pages)
    else:
        budget = pdf_ocr_page_cap(filename)
    batch = pdf_ocr_batch_pages()

    parts: list[str] = []
    meta: dict[str, Any] = {}
    ocr_used = 0
    ocr_attempts = 0
    empty_text_pages = 0
    for start in range(0, n_pages, batch):
        end = min(start + batch, n_pages)
        remaining = max(0, budget - ocr_attempts)
        (text, bmeta), diag = run_isolated(
            _extract_pdf_range_job,
            (file_path, filename, start, end, remaining),
            fallback=("", {}),
            label=f"extraction of {filename} pages {start + 1}-{end}",
        )
        if text:
            parts.append(text)
        ocr_used += int((bmeta or {}).get("ocr_pages") or 0)
        ocr_attempts += int((bmeta or {}).get("ocr_attempts") or 0)
        empty_text_pages += int((bmeta or {}).get("empty_text_pages") or 0)
        for key in (
            "ocr_low_quality",
            "ocr_truncated",
            "text_truncated",
            "extract_oom_partial",
            "extract_chars",
            "extract_parts",
        ):
            if (bmeta or {}).get(key):
                meta[key] = bmeta[key]
        if diag:
            meta.update(diag)
            meta["extract_partial"] = True
            break

    if ocr_skipped:
        meta["ocr_skipped_too_large"] = True
        meta["size_mb"] = round(size_mb, 1)
    if ocr_used:
        meta["ocr_pages"] = ocr_used
        meta["ocr_low_quality"] = True
    if ocr_attempts:
        meta["ocr_attempts"] = ocr_attempts
    if empty_text_pages:
        meta["empty_text_pages"] = empty_text_pages
    if empty_text_pages > 0 and ocr_attempts == 0:
        meta["ocr_required"] = True
    return "\n".join(parts), meta


def _extract_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint file, slide by slide.

    Collects text from all shapes on all slides, prefixed with the slide
    number so chunk context stays answerable. Never raises.
    """
    try:
        import pptx

        with file_crypto.open_plaintext(file_path) as readable_path:
            prs = pptx.Presentation(readable_path)
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"Slide {i}:\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception:
        return ""


def _extract_kmz(file_path: str) -> str:
    """Extract text from a KMZ file (ZIP containing KML).

    Reads the first .kml entry, strips XML tags, and returns the plain text.
    KMZ may contain images/models; we only index the KML narrative so the
    file is locatable by its name and any embedded labels. Never raises.
    """
    try:
        import re
        import zipfile

        with file_crypto.open_plaintext(file_path) as readable_path:
            with zipfile.ZipFile(readable_path, "r") as zf:
                kml_names = [n for n in zf.namelist() if n.lower().endswith(".kml")]
                if not kml_names:
                    return ""
                kml_bytes = zf.read(kml_names[0])
        kml_text = kml_bytes.decode("utf-8", errors="replace")
        # Metadata-only: index the human-meaningful <name>/<description> labels,
        # NOT the raw <coordinates> dump. A large KMZ otherwise explodes into
        # hundreds of coordinate-noise chunks (one file produced 810) that
        # pollute retrieval. Keeps the file locatable by its placemark names.
        labels = re.findall(r"<(?:name|description)>(.*?)</(?:name|description)>",
                            kml_text, flags=re.IGNORECASE | re.DOTALL)
        cleaned = []
        for lab in labels:
            lab = re.sub(r"<[^>]+>", " ", lab)  # strip any nested CDATA/markup
            lab = re.sub(r"\s+", " ", lab).strip()
            if lab and not re.fullmatch(r"[-\d.,\s]+", lab):  # drop pure coord strings
                cleaned.append(lab)
        # De-dup while preserving order; cap to keep the doc metadata-sized.
        seen, out = set(), []
        for lab in cleaned:
            if lab not in seen:
                seen.add(lab); out.append(lab)
        return " | ".join(out[:500])
    except Exception:
        return ""


def _extract_msg(file_path: str) -> str:
    """Extract text from an Outlook .msg file (OLE compound document).

    Reads the subject, sender, and plain-text body streams. This makes
    project correspondence searchable without relying on the full
    extract-msg package (which conflicts with our beautifulsoup4 pin).
    Never raises.
    """
    try:
        import olefile
    except Exception:
        return ""

    def _read_stream(ole, stream_name):
        if not ole.exists(stream_name):
            return None
        data = ole.openstream(stream_name).read()
        if stream_name.endswith("001F"):
            return data.decode("utf-16-le", errors="replace")
        return data.decode("cp1252", errors="replace")

    try:
        with file_crypto.open_plaintext(file_path) as readable_path:
            if not olefile.isOleFile(readable_path):
                return ""
            ole = olefile.OleFileIO(readable_path)
            try:
                subject = _read_stream(ole, "__substg1.0_0E04001F") or _read_stream(ole, "__substg1.0_0E04001E") or ""
                sender = _read_stream(ole, "__substg1.0_0C1A001F") or _read_stream(ole, "__substg1.0_0C1A001E") or ""
                body = _read_stream(ole, "__substg1.0_1000001F") or _read_stream(ole, "__substg1.0_1000001E") or ""
            finally:
                ole.close()
    except Exception:
        return ""

    parts = []
    if sender:
        parts.append(f"From: {sender}")
    if subject:
        parts.append(f"Subject: {subject}")
    if body:
        parts.append(body)
    return "\n\n".join(parts)


def _extract_doc(file_path: str) -> str:
    """Extract text from a legacy binary Microsoft Word .doc file.

    Tries external converters in order of reliability:
    ``antiword`` (Linux/Render), ``catdoc`` (Linux/Render), ``textract``
    (if installed), then Windows Word COM automation (if Word is installed
    and pywin32 is available). Never raises.
    """
    import shutil
    import subprocess

    converters = []
    antiword = shutil.which("antiword")
    catdoc = shutil.which("catdoc")
    if antiword:
        converters.append([antiword])
    if catdoc:
        converters.append([catdoc, "-d", "utf-8"])

    if not converters:
        # Neither converter on PATH. The remaining fallbacks cannot rescue a
        # Linux container -- textract is not a declared dependency and
        # win32com is Windows-only -- so every .doc silently extracted to ""
        # and indexed as ZERO_CHUNK. That was true in production until
        # antiword/catdoc were added to the image (2026-08-20); the failure
        # was invisible because the empty return looked like an empty file.
        # Say so once per call rather than letting the corpus fill with
        # contentless .doc rows.
        logger.warning(
            "no .doc converter on PATH (antiword/catdoc missing) — %s will "
            "extract to empty text; install antiword or catdoc in the image",
            os.path.basename(file_path),
        )

    for cmd in converters:
        try:
            with file_crypto.open_plaintext(file_path) as readable_path:
                result = subprocess.run(
                    cmd + [readable_path],
                    capture_output=True,
                    text=True,
                    errors="replace",
                    timeout=60,
                    env=scrubbed_env(),  # audit §6.1
                )
            text = (result.stdout or "").strip()
            if text:
                return text
        except Exception:
            continue

    # Optional textract fallback (not a declared dependency — user-install only).
    try:
        import textract

        with file_crypto.open_plaintext(file_path) as readable_path:
            return textract.process(readable_path).decode("utf-8", errors="replace").strip()
    except Exception:
        logger.warning(
            "swallowed %s in _extract_doc() — continuing",
            "Exception", exc_info=True,
        )

    # Optional Windows Word COM fallback (requires Word + pywin32).
    try:
        import win32com.client as win32

        with file_crypto.open_plaintext(file_path) as readable_path:
            word = win32.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(readable_path)
            try:
                return doc.Range().Text
            finally:
                doc.Close(SaveChanges=False)
                word.Quit()
    except Exception:
        return ""

    return ""


# Archive extraction guards (zip-bomb / runaway-nest protection).
_ARCHIVE_MAX_DEPTH = int(os.getenv("ARCHIVE_MAX_DEPTH", "3"))
_ARCHIVE_MAX_FILES = int(os.getenv("ARCHIVE_MAX_FILES", "100"))
_ARCHIVE_MAX_TOTAL_BYTES = int(os.getenv("ARCHIVE_MAX_TOTAL_BYTES", str(50 * 1024 * 1024)))
# Skip archives whose own compressed size exceeds this — extracting multi-GB
# archives on the 2 GB Render worker is a timeout/OOM risk.
_ARCHIVE_MAX_FILE_SIZE = int(os.getenv("ARCHIVE_MAX_FILE_SIZE", str(50 * 1024 * 1024)))


def _extract_archive(
    file_path: str,
    filename: str,
    opener,
    depth: int = 0,
    counters: dict[str, int] | None = None,
) -> str:
    """Recursively extract text from an archive (ZIP or RAR).

    ``opener`` must be a callable that takes a path and returns an object
    with ``namelist()`` and ``read(name)`` methods (``zipfile.ZipFile`` or
    ``rarfile.RarFile``).

    Guards against zip bombs and infinitely nested archives:
    * max depth ``_ARCHIVE_MAX_DEPTH``
    * max total files ``_ARCHIVE_MAX_FILES``
    * max total uncompressed bytes ``_ARCHIVE_MAX_TOTAL_BYTES``

    Nested archive contents are flattened into the same document's text,
    prefixed with their archive-internal path for provenance. Never raises.
    """
    if depth > _ARCHIVE_MAX_DEPTH:
        return ""
    if counters is None:
        counters = {"files": 0, "bytes": 0}
    try:
        if os.path.getsize(file_path) > _ARCHIVE_MAX_FILE_SIZE:
            return ""
    except Exception:
        return ""

    parts: list[str] = []
    try:
        with file_crypto.open_plaintext(file_path) as readable_path:
            with opener(readable_path) as archive:
                names = archive.namelist()
    except Exception:
        return ""

    for name in names:
        if counters["files"] >= _ARCHIVE_MAX_FILES:
            break
        lower = name.lower()
        ext = os.path.splitext(lower)[1]

        # Skip directories and macOS resource forks.
        if name.endswith("/") or "__macosx" in lower:
            continue

        try:
            with file_crypto.open_plaintext(file_path) as readable_path:
                with opener(readable_path) as archive:
                    info = archive.getinfo(name)
                    member_size = getattr(info, "file_size", getattr(info, "compress_size", 0)) or 0
                    if member_size > _ARCHIVE_MAX_TOTAL_BYTES:
                        continue
                    if counters["bytes"] + member_size > _ARCHIVE_MAX_TOTAL_BYTES:
                        break
                    data = archive.read(name)
        except Exception:
            continue

        counters["files"] += 1
        counters["bytes"] += len(data)
        if counters["bytes"] >= _ARCHIVE_MAX_TOTAL_BYTES:
            break

        # Write the member to a temp file so existing extractors can run.
        tmp_path: str | None = None
        try:
            fd, tmp_path = tempfile.mkstemp(suffix=ext or ".bin", prefix="fork_arc_")
            os.close(fd)
            with open(tmp_path, "wb") as fh:
                fh.write(data)

            # Nested archive: recurse.
            if ext in {".zip", ".rar"}:
                nested_opener = _zip_opener if ext == ".zip" else _rar_opener
                nested_text = _extract_archive(
                    tmp_path, name, nested_opener, depth=depth + 1, counters=counters
                )
                if nested_text:
                    parts.append(f"[archive:{name}]\n{nested_text}")
                continue

            # Skip images inside archives — per-photo OCR/YOLO is too expensive
            # when a ZIP contains hundreds of construction photos, and the archive
            # itself is still locatable by name plus any text/PDF members.
            if ext in _IMAGE_EXTS:
                continue

            # Supported file: route through the same extraction pipeline.
            if ext in _SUPPORTED_EXTS:
                member_text, _ = _extract_with_meta(tmp_path, name)
                if member_text:
                    parts.append(f"[file:{name}]\n{member_text}")
        except Exception:
            continue
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except OSError:
                    logger.warning(
                        "swallowed %s in _extract_archive() — continuing",
                        "OSError", exc_info=True,
                    )

    return "\n\n".join(parts)


def _zip_opener(path: str):
    import zipfile

    return zipfile.ZipFile(path, "r")


def _rar_opener(path: str):
    import rarfile

    return rarfile.RarFile(path)


def _extract_zip(file_path: str, filename: str) -> str:
    """Recursively extract text from a ZIP archive. Never raises."""
    return _extract_archive(file_path, filename, _zip_opener)


def _rar_available() -> bool:
    """Check whether rarfile can locate an unrar binary."""
    try:
        import rarfile
        rarfile.tool_setup()
        return bool(rarfile.UNRAR_TOOL)
    except Exception:
        return False


def _extract_rar(file_path: str, filename: str) -> str:
    """Recursively extract text from a RAR archive. Degrades to "" if no
    unrar binary is available (expected on Windows dev boxes). Never raises.
    """
    if not _rar_available():
        return ""
    return _extract_archive(file_path, filename, _rar_opener)


# ── helpers ───────────────────────────────────────────────────────────────────


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


def _data_dir() -> str:
    """Resolve DATA_DIR at call time (so tests can relocate it via setenv)."""
    data_dir = os.getenv("DATA_DIR", "./data")
    try:
        os.makedirs(data_dir, exist_ok=True)
    except OSError:
        data_dir = tempfile.gettempdir()
    return data_dir


def _ensure_sqlite_parent_dir() -> None:
    url = get_database_url()
    if url.startswith("sqlite:///"):
        parent = os.path.dirname(url[len("sqlite:///") :])
        if parent:
            os.makedirs(parent, exist_ok=True)


def _legacy_db_path() -> str:
    """Pre-unified layout: dedicated SQLite file under DATA_DIR."""
    return os.path.join(_data_dir(), "doc_index.db")


def _legacy_index_dir() -> str:
    """Pre-SQLite layout: one JSON file per project under data/doc_index/."""
    return os.path.join(_data_dir(), "doc_index")


# ── text extraction ───────────────────────────────────────────────────────────

def _extract_with_meta(
    file_path: str, filename: str, force_ocr: bool = False
) -> tuple[str, dict[str, Any]]:
    """Extract plaintext from a document, plus a small metadata dict.

    Returns ``(text, meta)`` where ``meta`` carries ``{"ocr_low_quality": True}``
    when the document was OCR'd and the OCR quality verdict flagged it as a poor
    scan (omitted otherwise). Never raises — returns ``("", {})`` on any error.

    Output is guaranteed NUL-free: malformed PDFs (broken font CID maps —
    the same files that spray MuPDF warnings) yield literal ``\\x00`` bytes in
    the text layer, and PostgreSQL rejects any INSERT whose text contains
    NUL, which zero-chunked whole documents during the P1b Drive ingestion.

    Supports:
    * .txt/.md/.csv/.json/.xml — via file_crypto.read_document
    * .pdf — fitz / PyMuPDF text layer; if that text is effectively empty
      (a scanned / image-only PDF) it falls back to OCR on the same file
    * .docx — python-docx; .xlsx — openpyxl
    * image extensions (.jpg/.png/.webp/...) — OCR via OCRBlock
    """
    # Extraction is the OOM path (see app/core/extract_isolated for the
    # measured curve), so it runs in a child process under RLIMIT_AS wherever
    # the platform allows. A runaway document then dies alone instead of
    # taking the single uvicorn worker -- and every request in flight -- down.
    from app.core.extract_isolated import run_isolated, worth_isolating

    if not worth_isolating(file_path):
        # Too small to threaten the box. Running in-process also keeps
        # extraction's side effects observable: a forked child returns only
        # (text, meta), so counters, caches and spies mutated during
        # extraction are lost with it.
        text, meta = _extract_with_meta_impl(file_path, filename, force_ocr)
        if "\x00" in text:
            text = text.replace("\x00", "")
        return text, meta

    _, ext = os.path.splitext((filename or "").lower())
    if ext == ".pdf" and _pdf_should_batch(file_path):
        # Large scans (the 20-32 MB priced-BOQ band): isolated page batches
        # so a later-range timeout keeps earlier OCR instead of ZERO_CHUNK.
        text, meta = _extract_pdf_batched(
            file_path, filename, force_ocr=force_ocr
        )
        if "\x00" in text:
            text = text.replace("\x00", "")
        return text, meta

    (text, meta), diag = run_isolated(
        _extract_with_meta_impl,
        (file_path, filename, force_ocr),
        fallback=("", {}),
        label=f"extraction of {filename}",
    )
    if diag:
        # Carry the reason, so a 0-chunk document says WHY rather than being
        # indistinguishable from a genuinely empty file.
        meta = {**(meta or {}), **diag}
    if "\x00" in text:
        text = text.replace("\x00", "")
    return text, meta


def _extract_with_meta_impl(
    file_path: str, filename: str, force_ocr: bool = False
) -> tuple[str, dict[str, Any]]:
    """Format dispatch for ``_extract_with_meta``. May return NUL bytes."""
    try:
        _, ext = os.path.splitext((filename or "").lower())
        if ext not in _SUPPORTED_EXTS:
            return "", {}

        # ── plain-text-like formats ──────────────────────────────────────────
        if ext in {".txt", ".md", ".csv", ".json", ".xml"}:
            raw = file_crypto.read_document(file_path)
            return raw.decode("utf-8", errors="replace"), {}

        # ── images → OCR + YOLO-World fallback ───────────────────────────────
        if ext in _IMAGE_EXTS:
            text, low_quality = _ocr_extract(file_path)
            # Construction photos usually OCR blank; run the baked YOLO-World
            # detector to produce a searchable summary of visible objects.
            yolo_text = _safety_world_extract(file_path, filename or "")
            if yolo_text:
                text = f"{text}\n\n{yolo_text}".strip() if text else yolo_text
            meta: dict[str, Any] = {}
            if low_quality:
                meta["ocr_low_quality"] = True
            return text, meta

        # ── PDF ──────────────────────────────────────────────────────────────
        # Per-page: keep each page's text layer, OCR image-only pages, add
        # table rows. Fixes the cover-page bug (a digital cover page no longer
        # suppresses OCR of image-only body pages). Memory-bounded.
        if ext == ".pdf":
            return _extract_pdf(file_path, filename, force_ocr=force_ocr)

        # ── DOC / DOCX ───────────────────────────────────────────────────────
        if ext == ".doc":
            return _extract_doc(file_path), {}
        if ext == ".docx":
            import docx
            with file_crypto.open_plaintext(file_path) as readable_path:
                document = docx.Document(readable_path)
                # python-docx `.paragraphs` EXCLUDES table cells. RFP/BOD/spec
                # docs put requirements, schedules and specs in tables — dropping
                # them under-extracts ~half the document (Anthropic/Kenya RFPs:
                # 5.6k para chars + 5.3k table chars → only ~2 chunks). Include
                # both.
                parts = [p.text for p in document.paragraphs if p.text.strip()]
                # getattr guard: real python-docx always exposes .tables, but a
                # malformed/partial document (or a test double) may not — don't
                # let a missing .tables drop the paragraph text we already have.
                for _tbl in getattr(document, "tables", None) or []:
                    for _row in _tbl.rows:
                        _cells = [c.text.strip() for c in _row.cells if c.text.strip()]
                        if _cells:
                            parts.append(" | ".join(_cells))
                return "\n".join(parts), {}

        # ── XLSX ─────────────────────────────────────────────────────────────
        if ext == ".xlsx":
            import openpyxl
            with file_crypto.open_plaintext(file_path) as readable_path:
                wb = openpyxl.load_workbook(readable_path, data_only=True)
                parts: list[str] = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    # Row-wise, not cell-wise: keep each row's cells together so
                    # a priced-BOQ line stays intact ("<desc> | <qty> | <unit> |
                    # <rate> | <total>"). Flattening every cell into one
                    # space-joined blob separates a rate from its item
                    # description and makes company unit rates unretrievable.
                    for row in ws:
                        cells = [
                            str(cell.value).strip()
                            for cell in row
                            if cell.value is not None and str(cell.value).strip()
                        ]
                        if cells:
                            parts.append(" | ".join(cells))
                return "\n".join(parts), {}

        # ── PPTX ─────────────────────────────────────────────────────────────
        if ext == ".pptx":
            return _extract_pptx(file_path), {}

        # ── KMZ ──────────────────────────────────────────────────────────────
        if ext == ".kmz":
            return _extract_kmz(file_path), {}

        # ── MSG (Outlook email) ──────────────────────────────────────────────
        if ext == ".msg":
            return _extract_msg(file_path), {}

        # ── ZIP / RAR ────────────────────────────────────────────────────────
        if ext == ".zip":
            return _extract_zip(file_path, filename), {}
        if ext == ".rar":
            return _extract_rar(file_path, filename), {}

    except MemoryError:
        # See _extract_pdf: a memory failure must never masquerade as an empty
        # document, or the caller indexes ZERO_CHUNK and reports nothing wrong.
        raise
    except Exception as exc:
        # Same rule for extractors that allocate through their own C libraries
        # and raise something other than MemoryError (PyMuPDF's RuntimeError
        # was the one that reached production). Promote it so the isolation
        # layer reports "ran out of memory" instead of "this file is empty".
        if _is_memory_exhaustion(exc):
            raise MemoryError(f"extraction exhausted memory: {exc}") from exc
        return "", {
            "extract_failed": type(exc).__name__,
            "extract_failed_detail": str(exc)[:300],
        }

    return "", {}


def extract_document_text(file_path: str, filename: str) -> str:
    """Extract plaintext from a document.

    Supports .txt/.md/.csv/.json/.xml (via file_crypto.read_document),
    .pdf (via fitz / PyMuPDF + open_plaintext, with OCR fallback for scanned
    image-only PDFs), .doc (antiword/catdoc/textract/Word COM), .docx (via
    python-docx + open_plaintext), .xlsx (via openpyxl + open_plaintext), .pptx
    (python-pptx), .kmz (KML text), .zip/.rar (recursive archive flattening),
    .msg (Outlook email body), and image formats (.jpg/.jpeg/.png/.webp/.gif/
    .bmp/.tif/.tiff) via OCR.

    Returns "" for unsupported extensions and on any extraction error —
    callers treat the empty string as "skipped". Never raises.
    """
    text, _ = _extract_with_meta(file_path, filename)
    return text


# ── chunking ──────────────────────────────────────────────────────────────────

def chunk_text(text: str, words_per_chunk: int = 500) -> list[str]:
    """Split ``text`` into chunks of at most ``words_per_chunk`` words.

    Splits on whitespace; drops empty/whitespace-only chunks. Returns [] for
    empty or whitespace-only input.
    """
    words = text.split()
    if not words:
        return []
    chunks: list[str] = []
    for i in range(0, len(words), words_per_chunk):
        chunk = " ".join(words[i : i + words_per_chunk])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


# ── BOQ-aware finer chunker (FOLLOW-UP #91) ────────────────────────────────
#
# Used for documents where word-count chunking produces too-coarse chunks
# (the the client BOQ symptom — 8 chunks at 2798 chars per chunk, so finding
# a single rate requires the LLM to scan 700 tokens). Re-chunks at character
# level with overlap and respects BOQ row boundaries (lines that start with
# an item code like ``D 999.14``) so a rate stays adjacent to its item code.

_BOQ_ITEM_PATTERN = re.compile(r"^\s*[A-Z]\s*\d+(?:\.\d+)+\b")

# ``re`` is used here; the module-level ``import re`` was added with the
# noise-filter helper earlier.


def chunk_text_with_overlap(
    text: str,
    target_chars: int = 500,
    overlap: int = 50,
    max_chars: int = 800,
) -> list[str]:
    """Split ``text`` into chunks targeting ``target_chars`` per chunk with
    ``overlap`` characters carried into the next chunk.

    Prefer breakpoints at, in order of preference:
    1. A line that looks like the start of a BOQ row (matches an item-code
       pattern such as ``D 999.14``) — keeps a row contiguous in one chunk.
    2. A paragraph break (``\\n\\n``).
    3. A line break (``\\n``).
    4. A sentence boundary (``. ``).
    5. Hard cut at ``max_chars`` if nothing better was found.

    ``overlap`` chars of trailing context from the previous chunk are
    prepended to the next so a row split mid-sentence still has context.

    Empty / whitespace-only input returns ``[]``.
    """

    if not text or not text.strip():
        return []
    if target_chars <= 0:
        raise ValueError("target_chars must be positive")
    if overlap < 0 or overlap >= target_chars:
        raise ValueError("overlap must be in [0, target_chars)")
    max_chars = max(max_chars, target_chars)

    chunks: list[str] = []
    i = 0
    n = len(text)
    while i < n:
        # Hard window — never go past max_chars from the start of this chunk.
        hard_end = min(i + max_chars, n)
        if hard_end == n:
            tail = text[i:].strip()
            if tail:
                chunks.append(tail)
            break

        # Soft target — prefer a break point near i+target_chars.
        soft_end = min(i + target_chars, n)
        search_start = max(soft_end - 120, i + target_chars // 2)

        # Look for BOQ row start within [soft_end, hard_end] — break ABOVE the
        # next row so the new chunk starts with that row's code.
        next_row = None
        for match in _BOQ_ITEM_PATTERN.finditer(text, soft_end, hard_end):
            # Walk back to the line start.
            line_start = text.rfind("\n", i, match.start()) + 1
            if line_start > i and line_start <= hard_end:
                next_row = line_start
                break
        if next_row is not None:
            end = next_row
        else:
            # Fall back to whitespace boundaries.
            end = -1
            for marker in ("\n\n", "\n", ". "):
                idx = text.rfind(marker, search_start, soft_end)
                if idx > 0:
                    end = idx + len(marker)
                    break
            if end <= i:
                # No boundary found in the soft window — extend to hard_end.
                end = hard_end

        chunk = text[i:end].strip()
        if chunk:
            chunks.append(chunk)
        # Carry overlap from the tail of this chunk.
        next_i = end - overlap if end - overlap > i else end
        if next_i <= i:
            next_i = end
        i = next_i
    return chunks


def chunk_markdown(
    text: str,
    target_chars: int = 500,
    overlap: int = 50,
    max_table_chars: int = 2200,
) -> list[str]:
    """Markdown-aware chunking for curated knowledge notes.

    Keeps each markdown TABLE atomic — header row + separator + all data rows,
    prefixed with the nearest preceding heading — as a single chunk, and
    char-chunks the prose between tables. Plain char-chunking split tables so a
    data row (e.g. "Silver | EPC | Contractor | Lump sum") ended up in a
    different chunk from its column header ("Book | Form | Design | Pricing
    basis"); retrieval then matched the headerless fragment poorly and the model
    couldn't map the value to its column. Keeping the table whole (with its
    heading) fixes both retrieval and extraction. A table larger than
    ``max_table_chars`` falls back to char-chunking so it can't blow the budget.
    """
    if not text or not text.strip():
        return []
    lines = text.split("\n")
    sep_re = re.compile(r"\s*\|[\s:|\-]+\|")
    segments: list[tuple] = []  # (kind, text)
    prose_buf: list[str] = []
    last_heading = ""

    def _flush_prose():
        if prose_buf:
            joined = "\n".join(prose_buf).strip()
            if joined:
                segments.append(("prose", joined))
            prose_buf.clear()

    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        stripped = line.strip()
        if stripped.startswith("#"):
            last_heading = stripped
        # A table = a "|" row immediately followed by a |---|---| separator.
        if (
            stripped.startswith("|")
            and i + 1 < n
            and sep_re.match(lines[i + 1])
        ):
            _flush_prose()
            tbl = [line]
            j = i + 1
            while j < n and lines[j].lstrip().startswith("|"):
                tbl.append(lines[j])
                j += 1
            body = "\n".join(tbl)
            table_text = f"{last_heading}\n{body}" if last_heading else body
            segments.append(("table", table_text))
            i = j
            continue
        prose_buf.append(line)
        i += 1
    _flush_prose()

    chunks: list[str] = []
    for kind, seg in segments:
        seg = seg.strip()
        if not seg:
            continue
        if kind == "table" and len(seg) <= max_table_chars:
            chunks.append(seg)
        else:
            chunks.extend(
                chunk_text_with_overlap(seg, target_chars=target_chars, overlap=overlap)
            )
    return chunks


def _chunker_for_document(filename: str, requested: str = "default") -> str:
    """Pick a chunker. BOQ-named files default to the finer item-code splitter.

    The 500-word window buried a single rate (D599.5) inside a wall of
    unrelated rows. Admin reindex can still pass ``chunker=finer`` explicitly;
    this makes the product path (upload / eager index / index_project) do
    the same for files whose names already say they are a BOQ.
    """
    if requested != "default":
        return requested
    if _BOQ_NAME_RE.search(filename or ""):
        return "finer"
    return "default"


def _chunk_prose_segment(text: str, chunker: str) -> list[str]:
    if not text or not text.strip():
        return []
    if chunker == "markdown":
        return chunk_markdown(text, target_chars=500, overlap=50)
    if chunker == "finer":
        return chunk_text_with_overlap(text, target_chars=500, overlap=50)
    return chunk_text(text)


def _normalize_cesmm_in_chunks(chunks: list[str]) -> list[str]:
    """Collapse OCR-spaced CESMM codes (``D 549.2`` → ``D549.2``) in chunks."""
    if not chunks:
        return chunks
    from app.core.rag.vector_store import normalize_cesmm_item_codes
    return [normalize_cesmm_item_codes(c) for c in chunks]


def chunk_extracted_document(
    text: str,
    chunker: str = "default",
    filename: str = "",
) -> list[str]:
    """Chunk extracted document text, keeping Contract Data rows intact.

    Prose (or the whole document when no Contract Data section is found)
    uses the requested ``chunker``. Contract Data / Appendix-to-Tender /
    Contract Particulars sections are *replaced* by row-preserving
    particulars chunks so the default 500-word splitter cannot cut a
    filled-in amount mid-row. Those particulars chunks are the only copy
    of that section in the index.

    CESMM item codes are space-collapsed before chunking so ``D 549.2``
    and ``D549.2`` share a token and a row stays adjacent to its figures.
    """
    if not text or not text.strip():
        return []
    from app.core.rag.vector_store import normalize_cesmm_item_codes
    text = normalize_cesmm_item_codes(text)
    from app.core.contract_data_chunks import (
        contract_data_particulars_chunks,
        contract_data_spans,
    )

    spans = contract_data_spans(text, filename)
    if not spans:
        return _chunk_prose_segment(text, chunker)
    chunks: list[str] = []
    prev = 0
    for start, end in spans:
        if start > prev:
            chunks.extend(_chunk_prose_segment(text[prev:start], chunker))
        cd_chunks = contract_data_particulars_chunks(text[start:end], filename)
        if cd_chunks:
            chunks.extend(cd_chunks)
        else:
            chunks.extend(_chunk_prose_segment(text[start:end], chunker))
        prev = end
    if prev < len(text):
        chunks.extend(_chunk_prose_segment(text[prev:], chunker))
    if not chunks:
        return _chunk_prose_segment(text, chunker)
    return chunks


# ── index persistence ─────────────────────────────────────────────────────────

def _index_from_row(row: DocIndex | None) -> dict[str, Any] | None:
    if row is None:
        return None
    data = row.index_json
    return data if isinstance(data, dict) else None


def _ensure_project_row(project_id: str, session: Session) -> None:
    """Satisfy doc_index.project_id FK when tests or legacy imports use bare ids."""
    if session.get(Project, project_id) is not None:
        return
    session.add(
        Project(
            id=project_id,
            name=project_id,
            client=None,
            status="active",
            aconex_connected=False,
            user_id="system",
            created_at=_now(),
        )
    )
    session.flush()


def _ensure_project_row_on_conn(conn: Connection, project_id: str) -> None:
    """FK helper for the SQLite BEGIN IMMEDIATE connection path."""
    exists = conn.execute(
        select(Project.id).where(Project.id == project_id)
    ).scalar_one_or_none()
    if exists is not None:
        return
    conn.execute(
        insert(Project).values(
            id=project_id,
            name=project_id,
            client=None,
            status="active",
            aconex_connected=False,
            user_id="system",
            created_at=_now(),
        )
    )


def _import_legacy_json_indexes(session: Session) -> None:
    """Import pre-SQLite data/doc_index/<pid>.json files once."""
    legacy = _legacy_index_dir()
    if not os.path.isdir(legacy):
        return
    for fn in os.listdir(legacy):
        if not fn.endswith(".json"):
            continue
        pid = fn[:-5]
        if _is_master_corpus_alias(pid):
            continue  # virtual alias — never materialise a real row for it
        if session.get(DocIndex, pid) is not None:
            continue
        try:
            with open(os.path.join(legacy, fn), "r", encoding="utf-8") as fh:
                data = json.load(fh)
        except (OSError, json.JSONDecodeError):
            continue
        if not isinstance(data, dict):
            continue
        _ensure_project_row(pid, session)
        session.add(
            DocIndex(project_id=pid, index_json=data, updated_at=_now())
        )


def _import_legacy_sqlite_db(session: Session) -> None:
    """Import rows from the legacy dedicated doc_index.db file once."""
    path = _legacy_db_path()
    if not os.path.isfile(path):
        return
    with sqlite3.connect(path, timeout=30.0) as leg:
        for pid, index_json, updated_at in leg.execute(
            "SELECT project_id, index_json, updated_at FROM doc_index"
        ):
            if _is_master_corpus_alias(pid):
                continue  # virtual alias — never materialise a real row for it
            if session.get(DocIndex, pid) is not None:
                continue
            try:
                data = json.loads(index_json)
            except (json.JSONDecodeError, TypeError):
                continue
            if not isinstance(data, dict):
                continue
            _ensure_project_row(pid, session)
            session.add(
                DocIndex(
                    project_id=pid,
                    index_json=data,
                    updated_at=updated_at or _now(),
                )
            )


def init_db() -> None:
    """Ensure the schema exists and import any legacy on-disk indexes.

    Pre-SQLite deployments stored each index as data/doc_index/<pid>.json.
    Pre-unified deployments used data/doc_index.db. Both are imported once
    per database URL (rows already present are left untouched); source files
    are left in place.
    """
    global _initialized, _initialized_for_url
    url = get_database_url()
    with _INDEX_LOCK:
        from app.core.projects import init_db as init_projects_db

        init_projects_db()
        _ensure_sqlite_parent_dir()
        DocIndex.__table__.create(bind=engine, checkfirst=True)
        if _initialized_for_url != url:
            with SessionLocal() as session:
                _import_legacy_json_indexes(session)
                # Flush pending DocIndex rows so the SQLite-migration's
                # ``session.get(DocIndex, pid)`` dedupe check can see them.
                # Without this, both migrations independently add a row for
                # any project_id present in both legacy sources, and the
                # final commit fails with the UNIQUE constraint.
                session.flush()
                _import_legacy_sqlite_db(session)
                session.commit()
            # Self-heal: drop any real row carrying the virtual master-corpus
            # alias id (a stale _ensure_project_row artifact that duplicates the
            # injected alias in listings).
            try:
                _purge_spurious_master_corpus_row()
            except Exception:
                _logging.getLogger(__name__).warning(
                    "init_db: spurious master-corpus row cleanup failed", exc_info=True
                )
            _initialized_for_url = url
        _initialized = True


def _ensure_db() -> None:
    url = get_database_url()
    if not _initialized or _initialized_for_url != url:
        init_db()


def _load_index(project_id: str) -> dict[str, Any] | None:
    """Read the stored index for ``project_id``. Returns None if absent."""
    _ensure_db()
    with SessionLocal() as session:
        return _index_from_row(session.get(DocIndex, project_id))


def _write_index(project_id: str, data: dict[str, Any]) -> None:
    """Replace the stored index for ``project_id`` (a full-rebuild write).

    Reuses ``_update_index`` so the full-rebuild path gets the same
    cross-process serialization as incremental updates:
    - SQLite: BEGIN IMMEDIATE
    - Postgres: pg_advisory_xact_lock + SELECT FOR UPDATE
    """
    _update_index(project_id, lambda _current: data)


def _apply_index_mutation(
    session: Session,
    project_id: str,
    mutate: Callable[[dict[str, Any] | None], dict[str, Any]],
    *,
    lock_row: bool,
) -> dict[str, Any]:
    if lock_row:
        row = session.scalar(
            select(DocIndex)
            .where(DocIndex.project_id == project_id)
            .with_for_update()
        )
    else:
        row = session.get(DocIndex, project_id)
    current_raw = _index_from_row(row)
    current = copy.deepcopy(current_raw) if current_raw is not None else None
    updated = mutate(current)
    now = _now()
    _ensure_project_row(project_id, session)
    if row is None:
        session.add(
            DocIndex(project_id=project_id, index_json=updated, updated_at=now)
        )
    else:
        row.index_json = updated
        row.updated_at = now
        flag_modified(row, "index_json")
    return updated


def _load_index_json_from_conn(
    conn: Connection, project_id: str
) -> dict[str, Any] | None:
    data = conn.execute(
        select(DocIndex.index_json).where(DocIndex.project_id == project_id)
    ).scalar_one_or_none()
    return data if isinstance(data, dict) else None


def _update_index_on_sqlite_conn(
    conn: Connection,
    project_id: str,
    mutate: Callable[[dict[str, Any] | None], dict[str, Any]],
) -> dict[str, Any]:
    current = _load_index_json_from_conn(conn, project_id)
    updated = mutate(current)
    now = _now()
    _ensure_project_row_on_conn(conn, project_id)
    row_exists = conn.execute(
        select(DocIndex.project_id).where(DocIndex.project_id == project_id)
    ).scalar_one_or_none()
    if row_exists is None:
        conn.execute(
            insert(DocIndex).values(
                project_id=project_id,
                index_json=updated,
                updated_at=now,
            )
        )
    else:
        conn.execute(
            update(DocIndex)
            .where(DocIndex.project_id == project_id)
            .values(index_json=updated, updated_at=now)
        )
    return updated


def _update_index(
    project_id: str,
    mutate: Callable[[dict[str, Any] | None], dict[str, Any]],
) -> dict[str, Any]:
    """Atomic read-modify-write of a project's index.

    Runs ``mutate(current_or_None) -> new_index`` inside a single write
    transaction (BEGIN IMMEDIATE on SQLite), so two concurrent updates —
    even from separate worker processes — serialise rather than overwrite
    each other.
    """
    with _INDEX_LOCK:
        _ensure_db()
        eng = get_engine()
        dialect = eng.dialect.name
        if dialect == "sqlite":
            with eng.connect() as conn:
                conn = conn.execution_options(isolation_level="AUTOCOMMIT")
                conn.exec_driver_sql("BEGIN IMMEDIATE")
                try:
                    updated = _update_index_on_sqlite_conn(conn, project_id, mutate)
                    conn.exec_driver_sql("COMMIT")
                except BaseException:
                    conn.exec_driver_sql("ROLLBACK")
                    raise
            return updated

        with SessionLocal() as session, session.begin():
            # FOR UPDATE does not lock missing rows; advisory lock serialises
            # concurrent first-time inserts for the same project_id.
            lock_key = zlib.crc32(project_id.encode("utf-8")) & 0x7FFFFFFF
            session.execute(
                text("SELECT pg_advisory_xact_lock(:k)"), {"k": lock_key}
            )
            return _apply_index_mutation(
                session, project_id, mutate, lock_row=True
            )


def _ext_of(filename: str) -> str:
    _, ext = os.path.splitext((filename or "").lower())
    return ext


def index_project(project_id: str) -> dict[str, Any]:
    """Build (or rebuild) the full text index for ``project_id``.

    For each document:
    - unsupported extension → record under ``skipped`` with reason
      ``"unsupported_type"``
    - supported extension → extract text → chunk → store entry

    Persists the result as JSON and returns a summary dict.
    """
    docs = _projects.list_documents(project_id)

    documents: list[dict[str, Any]] = []
    skipped: list[dict[str, Any]] = []
    total_chunks = 0

    for doc in docs:
        filename = doc.get("original_name", "")
        ext = _ext_of(filename)

        if ext not in _SUPPORTED_EXTS:
            fingerprint = f"{doc['uploaded_at']}:{doc['size']}"
            skipped.append({
                "document_id": doc["id"],
                "filename": filename,
                "reason": "unsupported_type",
                "fingerprint": fingerprint,
            })
            continue

        file_path = doc.get("file_path") or ""
        text, meta = _extract_with_meta(file_path, filename)
        chunks = chunk_extracted_document(
            text,
            chunker=_chunker_for_document(filename),
            filename=filename,
        )
        # BOQ → RAG: mirror index_document so the package total + line items
        # are retrievable regardless of which indexing path ran.
        boq_chunks = _boq_chunks_for_document(file_path, filename, ext, project_id)
        if boq_chunks:
            chunks = chunks + boq_chunks
        # Drawing → RAG: same mirroring rule as the BOQ hook. Both indexing
        # paths must append these or a schedule is retrievable after one kind
        # of ingest and invisible after the other.
        drawing_chunks = _drawing_chunks_for_document(file_path, filename, ext, project_id)
        if drawing_chunks:
            chunks = chunks + drawing_chunks
        ifc_chunks = _ifc_chunks_for_document(file_path, filename, ext, project_id)
        if ifc_chunks:
            chunks = chunks + ifc_chunks
        chunks = _normalize_cesmm_in_chunks(chunks)

        fingerprint = f"{doc['uploaded_at']}:{doc['size']}"
        entry: dict[str, Any] = {
            "document_id": doc["id"],
            "filename": filename,
            "fingerprint": fingerprint,
            "chunks": chunks,
        }
        if meta.get("ocr_low_quality"):
            entry["ocr_low_quality"] = True
        documents.append(entry)
        total_chunks += len(chunks)

        # PR #94: mirror the RAG hook from ``index_document`` so the
        # chunks table stays in sync with the JSON index. Without this,
        # ``index_project`` left the chunks table empty and the new
        # hybrid-retriever-backed ``search_project_documents`` returned
        # nothing — even though the JSON index was fresh.
        try:
            from app.core.rag import retriever as _rag
            if _rag.available() and chunks:
                _rag.index_chunks(project_id, doc["id"], chunks)
        except Exception as exc:  # noqa: BLE001 — RAG must never break primary indexing
            import logging as _logging
            _logging.getLogger(__name__).warning(
                "RAG indexing skipped for %s during index_project: %s",
                doc["id"], exc,
            )

    index_data: dict[str, Any] = {
        "project_id": project_id,
        "built_at": _now(),
        "documents": documents,
        "skipped": skipped,
    }
    _write_index(project_id, index_data)

    indexed = len(documents)
    skipped_unsupported = len(skipped)
    if indexed > 0 and total_chunks == 0:
        return {
            "status": "error",
            "error": "ZERO_CHUNK",
            "banner": (
                f"ERROR: re-index produced 0 chunks across {indexed} documents — "
                "check extractor/OCR logs."
            ),
            "project_id": project_id,
            "indexed": indexed,
            "skipped_unsupported": skipped_unsupported,
            "total_chunks": 0,
        }

    return {
        "status": "ok",
        "project_id": project_id,
        "indexed": indexed,
        "skipped_unsupported": skipped_unsupported,
        "total_chunks": total_chunks,
    }


# ── BOQ → RAG wiring ────────────────────────────────────────────────────────
# A bill-of-quantities document gets its COMPUTED total + line items appended
# as retrievable chunks, so "what is the total package value?" is answerable
# from the corpus. The raw spreadsheet/table text alone never contains the sum;
# boq_processor computes it. Best-effort: a parse failure never blocks the
# primary text indexing.

_BOQ_NAME_RE = re.compile(
    r"\bbo[q]\b|bill[\s_-]*of[\s_-]*quantit|schedule[\s_-]*of[\s_-]*quantit|priced",
    re.IGNORECASE,
)


def _looks_like_boq(filename: str, ext: str) -> bool:
    """Cheap gate before invoking the (heavier) BOQ parser.

    Spreadsheets are usually cost/quantity tables, so always attempt them
    (pandas parse is cheap). PDFs are attempted ONLY when the filename signals
    a BOQ, so we don't run pdfplumber table extraction on every PDF ingest.
    """
    if ext in {".xlsx", ".xls", ".csv"}:
        return True
    if ext == ".pdf" and _BOQ_NAME_RE.search(filename or ""):
        return True
    return False


def _boq_summary_chunks(result: dict[str, Any]) -> list[str]:
    """Stringify a boq_processor result into retrievable chunks.

    ACCURACY GUARD (no-assumptions rule): when ``pages_skipped > 0`` the total
    is reported as PARTIAL / INCOMPLETE — never a confident figure — so chat
    can never state a wrong package value with confidence.
    """
    if not result or result.get("status") != "success":
        return []
    line_items = result.get("line_items") or []
    total = result.get("total_cost")
    if not line_items and total is None:
        return []
    currency = (result.get("currency") or "").strip()
    pages_skipped = result.get("pages_skipped") or 0
    out: list[str] = []
    if total is not None:
        # Phrase the total the way users actually ask ("what is the total
        # package value?"), so this chunk out-matches generic contract
        # templates that merely mention "package value".
        if pages_skipped:
            out.append(
                "BOQ partial total — total package value / total contract value "
                "(some pages could not be parsed, this figure is INCOMPLETE): "
                f"{total} {currency}".strip()
            )
        else:
            out.append(
                "BOQ total — total package value, total contract value, total "
                f"cost of the bill of quantities: {total} {currency} "
                "(sum of all priced line items).".strip()
            )
    cost_breakdown = result.get("cost_breakdown") or {}
    if cost_breakdown:
        parts = [
            f"{section} {v.get('total')} {currency} ({v.get('percentage')}%)"
            for section, v in cost_breakdown.items()
        ]
        out.append("BOQ cost breakdown by section — " + "; ".join(parts))
    source = (result.get("source_name") or "").strip()
    src_prefix = f" [{source}]" if source else ""
    # Terminology alias so a "sewer" query matches "waste water" content and vice
    # versa (they are the SAME network; storm/surface water is a different one).
    _blob = " ".join(str(it.get("description") or "") for it in line_items).lower()
    if any(w in _blob for w in ("sewer", "foul", "waste water", "wastewater", "sanitary")):
        out.append(
            f"TERMINOLOGY{src_prefix}: in this Bill of Quantities the terms "
            "'sewer', 'foul sewer', 'foul drainage', 'sanitary drainage' and "
            "'waste water' all mean the SAME network. 'Storm water' / 'surface "
            "water' is a DIFFERENT, separate network. A question about 'sewer' "
            "quantities is answered by the 'waste water' / 'foul' items, and "
            "vice versa."
        )
    # NOTE: no software "aggregate total" chunk. Blindly summing OCR'd line-item
    # quantities is wrong — units are mixed and often mis-read (excavation /
    # backfill / concrete are m3, pipe-laying is LM, bedding is m3 or LM-by-
    # section), so a summed grand total is meaningless. Per-item chunks below
    # carry each measured quantity as-is; totals come from the operator's
    # verified figures, not from summing scanned line items.
    # Name the source BOQ on every line-item chunk so retrieval can disambiguate
    # WITHIN a project that holds several BOQs (e.g. the client project's three distinct
    # demolition totals must never be conflated). Phrase quantity/rate/total with
    # the synonyms operators actually type ("total quantity", "unit price",
    # "total price") plus the unit (UOM), so "unit price of X" / "total quantity
    # of X" retrieve cleanly.
    for item in line_items:
        desc = (item.get("description") or item.get("item_key") or "item").strip()
        unit = (item.get("unit") or "").strip()
        unit_str = f" {unit}" if unit else ""
        out.append(
            f"BOQ line item{src_prefix} — {desc}: total quantity "
            f"{item.get('quantity')}{unit_str}, unit price (rate) "
            f"{item.get('unit_cost')}, total price (amount) "
            f"{item.get('total_cost')} {currency}".strip()
        )
    return out


def _boq_no_total_guard(filename: str) -> str:
    """A guard chunk emitted when a document is clearly a BOQ but NO verified
    total could be computed from it (scanned/image PDF, or an unparseable
    table). It instructs the model NOT to state a total, so chat can't
    synthesise a confident-but-wrong 'total package value' from partial OCR
    text. Accuracy over coverage (no-assumptions rule)."""
    return (
        f"VERIFIED-TOTAL GUARD for '{filename}': this is a Bill of Quantities, "
        "but NO verified total could be computed from it — it appears to be a "
        "scanned/image PDF or its priced table could not be parsed, so only "
        "fragments were extracted. Do NOT state a total package value, total "
        "contract value, or any overall BOQ total from this document; any such "
        "figure would be partial or wrong. To get an accurate total, an "
        "xlsx/csv version of this BOQ is required."
    )


def _boq_chunks_for_document(
    file_path: str, filename: str, ext: str, project_id: str
) -> list[str]:
    """Return BOQ summary chunks for a BOQ-type document, or [] otherwise.

    When the file is clearly a BOQ (by filename) but no total could be
    computed — e.g. a scanned PDF pdfplumber can't parse — emit a guard chunk
    instead, so the model doesn't invent a total from partial OCR text.

    Never raises — a parse failure just yields the guard (or []) so the primary
    text indexing is unaffected.
    """
    if not file_path or not _looks_like_boq(filename, ext):
        return []
    try:
        from app.blocks.boq_processor import BOQProcessorBlock

        result = _run_sync(
            BOQProcessorBlock().process(
                {"file_path": file_path, "project_id": project_id}
            )
        )
        # Tag the result with the source BOQ name so per-line-item chunks can
        # name their origin (disambiguates several BOQs in one project).
        if isinstance(result, dict):
            result.setdefault("source_name", os.path.splitext(filename or "")[0])
        chunks = _boq_summary_chunks(result)
        if chunks:
            return chunks
        # No verified total. If the FILENAME clearly says BOQ (not just any
        # spreadsheet), emit the guard so chat refuses to state a false total.
        if _BOQ_NAME_RE.search(filename or ""):
            return [_boq_no_total_guard(filename)]
        return []
    except Exception as exc:  # noqa: BLE001 — never block indexing on a BOQ parse
        _logging.getLogger(__name__).warning(
            "BOQ wiring skipped for %s: %s", filename, exc
        )
        if _BOQ_NAME_RE.search(filename or ""):
            return [_boq_no_total_guard(filename)]
        return []


# ── Drawing tables → RAG ─────────────────────────────────────────────────────
#
# Same problem the BOQ hook above solves, one document type over. A drawing
# keeps its numbers in its SCHEDULES — bar/bending, door, window, finishes —
# and `_extract_with_meta` only pulls the PDF's raw text layer. Schedule cells
# come through as a flat run of words with the row/column structure gone, so
# "how many 16mm bars are on sheet S-2101-004" could not be answered even
# though the figure was on the page.
#
# `ConstructionContainer._extract_tables_advanced` already recovers the real
# table structure. This turns each recovered table into retrievable chunks,
# exactly as `_boq_chunks_for_document` does for priced line items.

_DRAWING_NAME_RE = re.compile(
    # `DWG` is the drawing token in the JCB/the client document-code convention
    # that names most of this corpus:
    #   IP-INF-053-0000-JCB-DWG-TM-200-1000005-A.pdf   <- drawing
    #   IP-INF-053-0000-JCB-BOQ-CA-000007-B_...pdf     <- NOT a drawing
    #   IP-INF-053-0000-JCB-SPC-IF-000013-B_SOPR.pdf   <- NOT a drawing
    # `\b` is WRONG here: underscore is a word character, so `\bdrawing\b`
    # does not match `drawing_tm_200.pdf` — a real filename from the corpus.
    # These lookarounds treat `_`, `-`, `.` and space alike as separators.
    #
    # Bare "sheet" is deliberately NOT a token. It is the loosest of the
    # candidates ("Data Sheet", "Method Statement Sheet") and the whole reason
    # this gate is filename-based is to keep table detection off PDFs that do
    # not need it. All four real drawing filenames in the corpus carry DWG or
    # "drawing", so "sheet" buys nothing and costs breadth.
    r"(?<![A-Za-z0-9])(?:dwgs?|drawings?)(?![A-Za-z0-9])",
    re.IGNORECASE,
)


def _looks_like_drawing(filename: str, ext: str) -> bool:
    """Cheap gate before opening a PDF and running table detection.

    Deliberately filename-only, and PDF-only, for the reason `_looks_like_boq`
    gives for its own gate: running table detection over every page of every
    PDF would add real cost to a bulk re-index, and re-encode runs on this
    corpus already sit close to a capacity wall. A drawing that is not named
    like one is missed; that is the accepted trade, and it is recorded in
    KNOWN_INCOMPLETE.md rather than left implicit.
    """
    return ext == ".pdf" and bool(_DRAWING_NAME_RE.search(filename or ""))


def _drawing_table_chunks(tables: list[dict[str, Any]], source: str) -> list[str]:
    """Serialise recovered drawing tables into retrievable chunks.

    One chunk per table, prefixed with the sheet it came from — a project
    holds many sheets and a bar schedule on one is not the bar schedule on
    another, the same disambiguation `_boq_summary_chunks` does with
    ``source_name``.

    The wording carries the synonyms operators actually type ("bar mark",
    "diameter", "quantity", "schedule"), for the same reason the BOQ line-item
    chunks do: a bare grid of numbers loses to generic boilerplate on cosine
    alone.
    """
    _KIND_WORDS = {
        "rebar_schedule": "reinforcement bar schedule / bar bending schedule (BBS) — "
                          "bar mark, diameter, length, shape code, quantity",
        "door_schedule": "door schedule — door reference, type, size, fire rating, quantity",
        "window_schedule": "window schedule — window reference, type, size, quantity",
        "finishes_schedule": "room finishes schedule — floor, wall and ceiling finishes",
        "revision_history": "drawing revision history — revision, date, description",
        "legend": "drawing legend / key — symbols and their meanings",
        "general_notes": "drawing general notes",
        "quantities": "quantities table — item, unit, quantity",
        "unclassified": "table",
    }
    out: list[str] = []
    src = f" [{source}]" if source else ""
    for table in tables:
        rows = table.get("rows") or []
        header = table.get("header") or []
        if not rows and not header:
            continue
        kind = table.get("kind", "unclassified")
        label = _KIND_WORDS.get(kind, _KIND_WORDS["unclassified"])
        lines = []
        if header:
            lines.append(" | ".join(header))
        for row in rows:
            lines.append(" | ".join(row))
        out.append(
            f"DRAWING TABLE{src} — {label}. "
            f"Sheet page {table.get('page')}, {table.get('row_count')} rows.\n"
            + "\n".join(lines)
        )
    return out


def _drawing_no_table_guard(filename: str) -> str:
    """Emitted when a document is named as a drawing but yields no readable
    schedule — a scanned/vector-only sheet whose text layer is empty.

    Mirrors `_boq_no_total_guard`, and for the same reason: without it, chat
    sees a drawing in the corpus, finds no schedule values, and is free to
    produce a plausible bar diameter or door count. This says: the sheet is
    here, its schedules are NOT readable, do not state values from it.
    """
    return (
        f"SCHEDULE-VALUE GUARD for '{filename}': this document is a drawing "
        "sheet, but NO schedule table could be read from it — it appears to be "
        "a scanned or image-only PDF with no recoverable text layer. Do NOT "
        "state bar marks, diameters, quantities, door or window references, or "
        "any other schedule value from this drawing; any such figure would be "
        "invented. To read its schedules, a vector (non-scanned) PDF of this "
        "sheet is required."
    )


def _drawing_chunks_for_document(
    file_path: str, filename: str, ext: str, project_id: str
) -> list[str]:
    """Return schedule + title-block chunks for a drawing, or [] otherwise.

    Never raises — a parse failure yields the guard (or []) so the primary
    text indexing is unaffected, exactly as `_boq_chunks_for_document` behaves.
    """
    if not file_path or not _looks_like_drawing(filename, ext):
        return []
    log = _logging.getLogger(__name__)
    try:
        import fitz  # PyMuPDF

        from app.containers.construction import ConstructionContainer

        container = ConstructionContainer()
        out: list[str] = []
        source = ""
        any_text = False
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                if page_num == 0:
                    # The title block is on sheet 1 and names the drawing, so
                    # every chunk below can be attributed to it.
                    raw = page.get_text() or ""
                    any_text = bool(raw.strip())
                    title_block = container._extract_title_block({"raw_text": raw})
                    source = title_block.get("drawing_number") or ""
                    if title_block.get("fields_found"):
                        named = ", ".join(
                            f"{k.replace('_', ' ')}: {v}"
                            for k, v in title_block.items()
                            if k != "fields_found" and v
                        )
                        out.append(
                            f"DRAWING TITLE BLOCK — {filename}. "
                            f"Sheet identity: {named}."
                        )
                elif not any_text and (page.get_text() or "").strip():
                    any_text = True
                out.extend(
                    _drawing_table_chunks(
                        container._extract_tables_advanced(page),
                        source or os.path.splitext(filename or "")[0],
                    )
                )
        if out:
            return out
        # Named like a drawing, opened fine, produced nothing readable. If
        # there was no text layer at all it is a scan — say so, so the model
        # cannot invent schedule values. If there WAS text but no tables, the
        # sheet genuinely has no schedule and the raw text already covers it.
        if not any_text:
            return [_drawing_no_table_guard(filename)]
        return []
    except Exception as exc:
        log.warning(
            "drawing table wiring skipped for %s: %s", filename, exc, exc_info=True
        )
        return []


_IFC_STEP_TYPES = (
    b"IFCWALLSTANDARDCASE",
    b"IFCWALL",
    b"IFCSLAB",
    b"IFCBEAM",
    b"IFCCOLUMN",
    b"IFCDOOR",
    b"IFCWINDOW",
    b"IFCSTAIR",
    b"IFCROOF",
    b"IFCSPACE",
    b"IFCBUILDINGSTOREY",
)


def _ifc_step_census_chunk(file_path: str, filename: str) -> list[str]:
    """Count IFC types from STEP bytes when the BIM extractor is unavailable."""
    try:
        raw = file_crypto.read_document(file_path)
    except Exception:
        try:
            with open(file_path, "rb") as fh:
                raw = fh.read(2_000_000)
        except Exception:
            return []
    blob = raw[:2_000_000].upper()
    counts: list[str] = []
    for token in _IFC_STEP_TYPES:
        n = blob.count(token)
        if n:
            counts.append(f"{token.decode('ascii')} {n}")
    if not counts:
        return []
    return [
        f"IFC MODEL — {filename}. STEP census: {', '.join(counts)}."
    ]


def _ifc_chunks_for_document(
    file_path: str, filename: str, ext: str, project_id: str
) -> list[str]:
    """Return a BIM census chunk for an IFC so it is not left 'Not indexed'.

    Leftover live UI: ``sample_office.ifc`` sat at chunk_count=0 because IFC
    was an unsupported extractor type. Clash stays off — leftover L2.
    Never raises.
    """
    del project_id  # signature matches the BOQ/drawing hooks
    if (ext or "").lower() != ".ifc" or not file_path:
        return []
    log = _logging.getLogger(__name__)
    try:
        from app.blocks.bim_extractor import BIMExtractorBlock

        result = _run_sync(
            BIMExtractorBlock().process(
                {"file_path": file_path},
                {"run_clash_detection": False},
            )
        )
        if isinstance(result, dict) and result.get("status") != "error":
            inner = (
                result.get("result")
                if isinstance(result.get("result"), dict)
                else result
            )
            schema = inner.get("ifc_schema") or inner.get("schema") or "IFC"
            qty = inner.get("quantities") or {}
            counts = []
            for cat, body in qty.items():
                n = body.get("count") if isinstance(body, dict) else body
                if isinstance(n, (int, float)) and n:
                    counts.append(f"{cat} {int(n)}")
            census = ", ".join(counts) if counts else "no counted categories"
            return [
                f"IFC MODEL — {filename}. Schema: {schema}. "
                f"Element census: {census}."
            ]
    except Exception as exc:  # noqa: BLE001
        log.warning("IFC census wiring skipped for %s: %s", filename, exc, exc_info=True)
    return _ifc_step_census_chunk(file_path, filename)


def _scanned_pdf_missing_ocr(ext: str, meta: dict[str, Any]) -> bool:
    """True when a PDF had empty body pages and OCR was never invoked.

    Live reindex of 20ac033d returned status=ok with 6 cover-page chunks
    because the finer chunker recast the text layer + VERIFIED-TOTAL GUARD
    as success. Item codes D599.5 / D549.2 were never in the extract.

    OCR that ran and returned empty (a genuinely blank digital page) is
    not this failure — ``ocr_attempts > 0`` means the trigger fired.
    """
    if (ext or "").lower() != ".pdf":
        return False
    if int(meta.get("ocr_pages") or 0) > 0:
        return False
    if int(meta.get("ocr_attempts") or 0) > 0:
        return False
    if meta.get("ocr_required"):
        return True
    if meta.get("ocr_skipped_too_large"):
        return True
    if int(meta.get("empty_text_pages") or 0) > 0:
        return True
    return False


def index_document(
    project_id: str,
    document_id: str,
    chunker: str = "default",
    *,
    force_ocr: bool = False,
) -> dict[str, Any]:
    """Incrementally index a single document into the project's index.

    Loads the existing index (or starts an empty one), extracts + chunks the
    given document, replaces any existing entry for that document_id, and
    writes back.

    Returns a summary dict with ``indexed`` count (always 1 on success).

    ``chunker`` selects the chunking strategy:

    * ``"default"`` — 500-word windows (legacy, what every doc uses today).
    * ``"finer"``   — char-level windows with overlap (FOLLOW-UP #91 fix
      for the the client BOQ symptom: coarse word-windows produced 8 chunks
      averaging 2800 chars each, so single line items were buried inside
      a wall of unrelated text).

    ``force_ocr`` (admin reindex) OCRs empty-text pages even when a leftover
    size gate would skip them, and uses the 400-page force budget so a
    370-page scanned BOQ can land item codes.
    """
    doc = _projects.get_document(document_id)
    if doc is None:
        return {"project_id": project_id, "indexed": 0, "error": "document not found"}

    filename = doc.get("original_name", "")
    ext = _ext_of(filename)
    fingerprint = f"{doc.get('uploaded_at', '')}:{doc.get('size', 0)}"

    # Slow work (text extraction, OCR, chunking) runs OUTSIDE the lock so it
    # does not serialise all indexing — only the load-modify-write below does.
    entry: dict[str, Any] | None = None
    skipped_entry: dict[str, Any] | None = None
    chunks: list[str] = []
    if ext not in _SUPPORTED_EXTS:
        skipped_entry = {
            "document_id": document_id,
            "filename": filename,
            "reason": "unsupported_type",
            "fingerprint": fingerprint,
        }
    else:
        file_path = doc.get("file_path") or ""
        text, meta = _extract_with_meta(
            file_path, filename, force_ocr=force_ocr
        )
        chunks = chunk_extracted_document(
            text,
            chunker=_chunker_for_document(filename, chunker),
            filename=filename,
        )
        # BOQ → RAG: append the computed total + line items so the package
        # value is answerable from the corpus (the raw text never holds the sum).
        boq_chunks = _boq_chunks_for_document(file_path, filename, ext, project_id)
        if boq_chunks:
            chunks = chunks + boq_chunks
        # Drawing → RAG: mirrors index_project above. A drawing's schedules
        # live in its tables, and the raw text layer loses their structure.
        drawing_chunks = _drawing_chunks_for_document(file_path, filename, ext, project_id)
        if drawing_chunks:
            chunks = chunks + drawing_chunks
        ifc_chunks = _ifc_chunks_for_document(file_path, filename, ext, project_id)
        if ifc_chunks:
            chunks = chunks + ifc_chunks
        chunks = _normalize_cesmm_in_chunks(chunks)
        if _scanned_pdf_missing_ocr(ext, meta):
            # Cover-page text + VERIFIED-TOTAL GUARD still produce chunks.
            # That must not look like a successful body index.
            import logging as _logging
            _logging.getLogger(__name__).error(
                "OCR_REQUIRED project=%s doc=%s file=%s empty_text_pages=%s "
                "ocr_pages=%s ocr_skipped=%s — scanned PDF indexed without "
                "body OCR; item codes will not be searchable",
                project_id, document_id, filename,
                meta.get("empty_text_pages", 0),
                meta.get("ocr_pages", 0),
                bool(meta.get("ocr_skipped_too_large")),
            )
            return {
                "status": "error",
                "error": "OCR_REQUIRED",
                "banner": (
                    "ERROR: scanned PDF produced cover-page chunks without "
                    "OCR — body item codes are missing. Reindex with "
                    "force_ocr=true (POST /v1/admin/debug/doc-reindex)."
                ),
                "project_id": project_id,
                "document_id": document_id,
                "filename": filename,
                "indexed": 0,
                "skipped_unsupported": 0,
                "total_chunks": len(chunks),
                "ocr_pages": int(meta.get("ocr_pages") or 0),
                "ocr_attempts": int(meta.get("ocr_attempts") or 0),
                "empty_text_pages": int(meta.get("empty_text_pages") or 0),
                "ocr_skipped_too_large": bool(meta.get("ocr_skipped_too_large")),
            }
        if not chunks:
            # A supported file that produced zero chunks is a SILENT indexing
            # failure: the document "exists" but can never be retrieved. Three
            # whole projects sat in that state for weeks before the 2026-06
            # audit noticed. WARNING (not error) because empty-but-valid files
            # exist; the marker makes the condition grep-able in Render logs.
            extract_error = None
            if meta.get("extract_failed"):
                detail = meta.get("extract_failed_detail") or ""
                extract_error = f"{meta['extract_failed']}: {detail}".rstrip(": ")
            import logging as _logging
            _logging.getLogger(__name__).warning(
                "ZERO_CHUNK project=%s doc=%s file=%s ext=%s — supported type "
                "extracted no chunks (empty file, failed OCR, or extractor bug)%s",
                project_id, document_id, filename, ext,
                f" extract_error={extract_error}" if extract_error else "",
            )
            result = {
                "status": "error",
                "error": "ZERO_CHUNK",
                "banner": (
                    "ERROR: document produced 0 chunks — check extractor/OCR."
                ),
                "project_id": project_id,
                "document_id": document_id,
                "filename": filename,
                "indexed": 0,
                "skipped_unsupported": 0,
                "total_chunks": 0,
            }
            if extract_error:
                result["extract_error"] = extract_error
            return result
        entry = {
            "document_id": document_id,
            "filename": filename,
            "fingerprint": fingerprint,
            "chunks": chunks,
        }
        if meta.get("ocr_low_quality"):
            entry["ocr_low_quality"] = True

        # ── RAG hook (PR 2) — best-effort embedding into the vector store.
        # Lazy import + try/except keep doc_index importable even when
        # sentence-transformers isn't installed. Idempotent via
        # upsert_chunks: re-indexing the same doc replaces its chunks.
        rag_indexed = 0
        rag_error: str | None = None
        try:
            from app.core.rag import retriever as _rag
            if not _rag.available():
                rag_error = "embedding stack unavailable"
            elif chunks:
                rag_indexed = _rag.index_chunks(project_id, document_id, chunks) or 0
                if rag_indexed:
                    entry["rag_indexed"] = rag_indexed
                else:
                    # Text was extracted but nothing reached the vector store,
                    # so the document is unretrievable by semantic search. Same
                    # end state as ZERO_CHUNK, different cause — and it used to
                    # look identical to success.
                    rag_error = f"embedded 0 of {len(chunks)} chunks"
        except Exception as exc:  # noqa: BLE001
            # Never let a RAG failure abort the primary doc-index path
            rag_error = f"{type(exc).__name__}: {exc}"
            import logging as _logging
            import traceback as _traceback
            _logging.getLogger(__name__).warning(
                "RAG indexing traceback for %s:\n%s", document_id, _traceback.format_exc()
            )

        if rag_error:
            # RAG_INDEX_FAILED is the grep-able marker, and this is an ERROR,
            # not a warning. A document that extracts text but never reaches the
            # vector store is stored, registered and listed in the sidebar while
            # being permanently unsearchable — the exact silent state the
            # ZERO_CHUNK marker above exists to prevent, reached by a different
            # route. It was logged as a warning and dropped, so a failed
            # embedder download (e.g. a 403 from the model host) presented to
            # the user as a successful upload that simply never had answers in
            # it. Recorded on the entry as well so the condition is queryable
            # after the fact, not only findable in log scrollback.
            import logging as _logging
            entry["rag_error"] = rag_error
            _logging.getLogger(__name__).error(
                "RAG_INDEX_FAILED project=%s doc=%s file=%s chunks=%d — %s; "
                "document is stored but NOT semantically searchable",
                project_id, document_id, filename, len(chunks), rag_error,
            )

    # Load-modify-write inside one SQLite transaction — a concurrent
    # index_document call for the same project (another BackgroundTask, or
    # another worker process) cannot interleave and drop this entry.
    def _mutate(current: dict[str, Any] | None) -> dict[str, Any]:
        current = current or {
            "project_id": project_id,
            "built_at": _now(),
            "documents": [],
            "skipped": [],
        }
        current["documents"] = [
            d for d in current.get("documents", [])
            if d["document_id"] != document_id
        ]
        current["skipped"] = [
            s for s in current.get("skipped", [])
            if s["document_id"] != document_id
        ]
        if entry is not None:
            current["documents"].append(entry)
        else:
            current["skipped"].append(skipped_entry)
        current["built_at"] = _now()
        return current

    _update_index(project_id, _mutate)

    if entry is None:
        return {
            "status": "ok",
            "project_id": project_id,
            "indexed": 0,
            "skipped_unsupported": 1,
            "total_chunks": 0,
        }
    result = {
        "status": "ok",
        "project_id": project_id,
        "indexed": 1,
        "skipped_unsupported": 0,
        "total_chunks": len(chunks),
        "rag_indexed": rag_indexed,
        "ocr_pages": int(meta.get("ocr_pages") or 0),
        "ocr_attempts": int(meta.get("ocr_attempts") or 0),
        "empty_text_pages": int(meta.get("empty_text_pages") or 0),
    }
    if rag_error:
        # Surfaced to the CALLER, not just the log. The upload path reports
        # indexing status back to the user; without this it reported success
        # for a document that cannot be retrieved.
        result["rag_error"] = rag_error
    return result


def invalidate_project(project_id: str) -> None:
    """Delete the stored index for ``project_id`` (next access rebuilds it)."""
    with _INDEX_LOCK:
        _ensure_db()
        with SessionLocal() as session:
            session.execute(
                delete(DocIndex).where(DocIndex.project_id == project_id)
            )
            session.commit()


def _is_master_corpus_alias(project_id: str) -> bool:
    """True when ``project_id`` is the VIRTUAL master-corpus alias (backed by a
    different source project). Such an id must never become a real Project row —
    the alias is injected into listings on the fly."""
    from app.core.projects import (
        MASTER_CORPUS_PROJECT_ID,
        MASTER_CORPUS_SOURCE_PROJECT_ID,
    )
    return (
        project_id == MASTER_CORPUS_PROJECT_ID
        and MASTER_CORPUS_PROJECT_ID != MASTER_CORPUS_SOURCE_PROJECT_ID
    )


def purge_project_index(project_id: str) -> None:
    """Remove ALL index traces of a project so a delete sticks across restarts.

    ``delete_project`` only removes the Project row (+ cascaded docs). The
    DocIndex row and the legacy on-disk sources (``data/doc_index/<pid>.json``
    and the legacy sqlite db) survive — and ``init_db`` re-imports them on the
    next restart, where ``_ensure_project_row`` resurrects the deleted project.
    Purging those sources is what makes a delete permanent.
    """
    log = _logging.getLogger(__name__)
    with _INDEX_LOCK:
        _ensure_db()
        with SessionLocal() as session:
            session.execute(
                delete(DocIndex).where(DocIndex.project_id == project_id)
            )
            session.commit()
    # Legacy per-project JSON file.
    try:
        path = os.path.join(_legacy_index_dir(), f"{project_id}.json")
        if os.path.isfile(path):
            os.remove(path)
    except OSError as exc:
        log.warning("purge_project_index: legacy json remove failed for %s: %s", project_id, exc)
    # Legacy dedicated sqlite db row.
    try:
        db = _legacy_db_path()
        if os.path.isfile(db):
            with sqlite3.connect(db, timeout=30.0) as leg:
                leg.execute("DELETE FROM doc_index WHERE project_id = ?", (project_id,))
                leg.commit()
    except sqlite3.Error as exc:
        log.warning("purge_project_index: legacy db purge failed for %s: %s", project_id, exc)


def _purge_spurious_master_corpus_row() -> None:
    """Remove a real Project row that carries the virtual master-corpus alias id.

    The master corpus is injected into listings from its backing source project;
    a real row with the alias id is always a stale ``_ensure_project_row``
    artifact that duplicates the alias in the UI. Self-healing: runs on startup.
    """
    from app.core.projects import (
        MASTER_CORPUS_PROJECT_ID,
        MASTER_CORPUS_SOURCE_PROJECT_ID,
    )
    if MASTER_CORPUS_PROJECT_ID == MASTER_CORPUS_SOURCE_PROJECT_ID:
        return  # aliasing disabled — the id IS a real project
    with SessionLocal() as session:
        row = session.get(Project, MASTER_CORPUS_PROJECT_ID)
        if row is not None:
            session.delete(row)
            session.commit()
    # Delete the matching DocIndex row directly; do NOT call purge_project_index
    # because that calls _ensure_db() and can recurse back into init_db while
    # tables are still being created. A spurious alias row never has legacy files.
    with SessionLocal() as session:
        session.execute(
            delete(DocIndex).where(DocIndex.project_id == MASTER_CORPUS_PROJECT_ID)
        )
        session.commit()


# ── search ────────────────────────────────────────────────────────────────────

async def search_project_documents(
    project_id: str,
    query: str,
    top_k: int = 5,
) -> list[dict[str, Any]]:
    """Search indexed documents for ``query``, returning up to ``top_k`` results.

    Each result is a dict: ``{document_id, filename, snippet, score}``.

    PR #94: routes through the production hybrid retriever
    (``app.core.rag.retriever.retrieve_with_filter`` — BM25 + vector RRF
    over the ``chunks`` Postgres table) so the agent's tool sees the
    SAME chunks the RAG injection layer sees. Pre-PR-94 this used a
    separate TF-IDF-over-JSON-blobs path that did not query the migrated
    drive_archive corpus, causing tool results to be empty while the
    injected RAG context contained the right answer (PR #93 migration
    surfaced the gap).

    Best chunk per document is kept (so a ``top_k=5`` request returns up
    to 5 distinct files, not 5 chunks from one file). Filenames are
    resolved via ``_doc_name_for_id`` (the same helper the retriever
    itself uses).
    """
    if not query or not query.strip():
        return []

    # Hybrid retriever lives in app.core.rag.retriever; local import to
    # keep app.core.doc_index importable when the RAG stack isn't fully
    # wired (tests, minimal install).
    from app.core.rag.retriever import _doc_name_for_id, retrieve_with_filter

    # Over-fetch chunks so we can still return ``top_k`` DISTINCT documents
    # after the "best chunk per document" collapse below. 4x is enough
    # for typical corpora where each doc has 2-10 chunks.
    over_fetch = max(top_k * 4, 20)

    def _query() -> list[Any]:
        try:
            chunks, _noise = retrieve_with_filter(query, project_id, k=over_fetch)
            return chunks
        except Exception:
            return []

    chunks = _query()
    # Lazy bootstrap: pre-PR-94 callers relied on the legacy code path
    # to build the project index on first search. Preserve that contract
    # for newly-uploaded projects (or anything pre-RAG-migration) by
    # iterating documents through ``index_document`` once. That call (and
    # not ``index_project``) is the one that fires ``_rag.index_chunks``,
    # which populates the ``chunks`` table the hybrid retriever queries.
    if not chunks and _load_index(project_id) is None:
        try:
            for doc in _projects.list_documents(project_id):
                try:
                    index_document(project_id, doc["id"])
                except Exception:
                    # Best-effort per-doc indexing; one failure shouldn't
                    # abort the bootstrap for the rest of the project.
                    continue
        except Exception:
            logger.warning(
                "swallowed %s in search_project_documents() — continuing",
                "Exception", exc_info=True,
            )
        chunks = _query()
    if not chunks:
        return []

    best: dict[str, dict[str, Any]] = {}
    for c in chunks:
        score = float(c.score or 0.0)
        prev = best.get(c.doc_id)
        if prev is None or score > prev["score"]:
            best[c.doc_id] = {
                "document_id": c.doc_id,
                "filename": _doc_name_for_id(c.doc_id),
                "chunk": c.text,
                "score": score,
                # Where this hit actually came from: "own" | "general_knowledge"
                # | "master_corpus". `Chunk.layer` is withheld from the LLM path
                # by design (the chat runtime reads it to phrase its own
                # disclosure). This is a DIRECT API consumer, not model context
                # — and without it a caller cannot tell an own-document hit from
                # a STEP 0b Master-Corpus fallback hit. Verified live
                # 2026-08-02: a project owning ZERO documents returned five
                # Master-Corpus documents here with nothing marking them.
                "origin": getattr(c, "layer", "own") or "own",
            }

    ranked = sorted(best.values(), key=lambda x: x["score"], reverse=True)[:top_k]

    # Pull the OCR-low-quality flag from the JSON index (if present). The
    # hybrid retriever's Chunk doesn't carry it; tests + UI need it for
    # the "low-confidence OCR" badge on citations.
    legacy_index = _load_index(project_id)
    ocr_flag_by_doc: dict[str, bool] = {}
    if legacy_index is not None:
        for entry in legacy_index.get("documents", []):
            if entry.get("ocr_low_quality"):
                ocr_flag_by_doc[entry["document_id"]] = True

    results: list[dict[str, Any]] = []
    for item in ranked:
        snippet = " ".join(item["chunk"].split()[:50])
        result_row = {
            "document_id": item["document_id"],
            "filename": item["filename"],
            "snippet": snippet,
            "score": round(item["score"], 4),
            # Which corpus this hit came from. This row is rebuilt from
            # scratch rather than passed through, so anything added to the
            # `best` dict above must ALSO be copied here or it never reaches
            # the caller — that is exactly how the first attempt at this
            # field shipped as a silent no-op.
            "origin": item.get("origin", "own"),
        }
        if ocr_flag_by_doc.get(item["document_id"]):
            result_row["ocr_low_quality"] = True
        results.append(result_row)
    return results


# ── eager (background) indexing ───────────────────────────────────────────────

def maybe_eager_index(project_id: str, document_id: str) -> dict[str, Any] | None:
    """Index a single document if eager indexing is enabled (default: on).

    Checks INDEX_ON_UPLOAD env var at call time — "1", "true", or "yes"
    (case-insensitive) enables; anything else disables. Intended to be
    scheduled via FastAPI BackgroundTasks so it runs after the response
    is sent without blocking the upload.

    BackgroundTasks discard the return value, so indexing status is also
    persisted onto the document row (``metadata.indexing``) for listing.
    """
    if os.getenv("INDEX_ON_UPLOAD", "true").strip().lower() not in ("1", "true", "yes"):
        return None
    result = index_document(project_id, document_id)
    try:
        indexing = {
            "status": result.get("status") or ("error" if result.get("error") else "ok"),
            "error": result.get("error"),
            "chunks": result.get("total_chunks", 0),
            "detail": result.get("extract_error") or result.get("banner") or result.get("rag_error"),
        }
        _projects.update_document_metadata(document_id, {"indexing": indexing})
    except Exception:
        logger.warning(
            "failed to persist indexing status for %s/%s",
            project_id, document_id, exc_info=True,
        )
    return result
